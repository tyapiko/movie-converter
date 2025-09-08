import streamlit as st
import tempfile
import os
from moviepy import VideoFileClip, AudioFileClip, CompositeAudioClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont
import numpy as np
import cv2
from pptx import Presentation
import io

# ✅ 実験完了: GitHub Actionsが構文エラーを正常に検出しました

st.set_page_config(
    page_title="動画編集ツール",
    page_icon="🎬",
    layout="wide"
)

# サイドバーでツール選択
st.sidebar.title("🛠️ ツール選択")
tool = st.sidebar.radio(
    "使用するツール",
    ["ショート動画変換", "動画結合", "パワポナレーション動画"]
)

if tool == "ショート動画変換":
    st.title("🎬 ショート動画コンバーター8")
    st.markdown("動画をアップロードして、YouTubeショート向けの縦型動画に変換しましょう！")
elif tool == "動画結合":
    st.title("🔗 動画結合ツール")
    st.markdown("複数のショート動画を選択して、長時間動画に結合しましょう！")
else:  # パワポナレーション動画
    st.title("📊 パワポナレーション動画")
    st.markdown("PowerPointファイルをアップロードして、ノート部分を読み上げる動画を作成しましょう！")

def resize_video_to_shorts(video_path, output_path, scale_factor=1.0, start_time=None, end_time=None, keep_original_size=False):
    """動画をYouTubeショート形式(9:16)にリサイズ、または元のサイズを維持"""
    import subprocess
    
    # FFmpegコマンドで動画変換
    ffmpeg_cmd = ['ffmpeg', '-i', video_path]
    
    # トリミングが指定されている場合
    if start_time is not None and end_time is not None:
        ffmpeg_cmd.extend(['-ss', str(start_time), '-t', str(end_time - start_time)])

    if not keep_original_size:
        # 元の動画情報を取得
        clip = VideoFileClip(video_path)
        original_width, original_height = clip.size
        original_ratio = original_width / original_height
        clip.close()
        
        # YouTubeショートの推奨解像度: 1080x1920 (9:16)
        target_width = 1080
        target_height = 1920
        target_ratio = target_width / target_height
        
        # 基本スケール計算（ターゲット枠に収まるサイズ）
        if original_ratio > target_ratio:
            # 横長の場合、幅をターゲット幅に合わせる
            base_width = target_width
            base_height = int(target_width / original_ratio)
        else:
            # 縦長の場合、高さをターゲット高さに合わせる
            base_height = target_height
            base_width = int(target_height * original_ratio)
        
        # scale_factorを適用（拡大倍率による調整）
        final_width = int(base_width * scale_factor)
        final_height = int(base_height * scale_factor)
        
        # 拡大倍率が1.0より大きい場合、動画がターゲットフレームからはみ出すのは正常
        # パディングエラーを避けるため、最小サイズは1ピクセル以上を保証
        final_width = max(1, final_width)
        final_height = max(1, final_height)
        
        # ビデオフィルターを構築
        if scale_factor > 1.0:
            # 拡大時：スケール→中央クロップ→パディング
            vf = f'scale={final_width}:{final_height},crop={min(final_width, target_width)}:{min(final_height, target_height)},pad={target_width}:{target_height}:(ow-iw)/2:(oh-ih)/2:black'
        else:
            # 縮小時：スケール→パディング
            vf = f'scale={final_width}:{final_height},pad={target_width}:{target_height}:(ow-iw)/2:(oh-ih)/2:black'
        
        ffmpeg_cmd.extend([
            '-vf', vf
        ])
    
    ffmpeg_cmd.extend([
        '-c:v', 'libx264',
        '-c:a', 'aac',
        '-b:v', '8000k',
        '-crf', '18',
        '-preset', 'slow',
        '-y',  # overwrite output file
        output_path
    ])
    
    try:
        subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
        return output_path
    except subprocess.CalledProcessError as e:
        raise Exception(f"FFmpeg処理でエラーが発生しました: {e.stderr.decode()}")
    except FileNotFoundError:
        raise Exception("FFmpegが見つかりません。システムにFFmpegがインストールされていることを確認してください。")

def get_voicevox_url():
    """VOICEVOX接続URLを取得（環境変数を優先）"""
    import os
    
    # 環境変数から取得（Docker環境で設定）
    voicevox_url = os.getenv('VOICEVOX_URL')
    if voicevox_url:
        return voicevox_url
    
    # Docker Compose環境では voicevox サービス名で接続
    try:
        # Docker環境かチェック
        if os.path.exists('/.dockerenv'):
            return "http://voicevox:50021"
    except:
        pass
    
    # WSL環境の場合はWindowsホストIPを取得
    try:
        with open('/etc/resolv.conf', 'r') as f:
            for line in f:
                if line.startswith('nameserver'):
                    host_ip = line.split()[1]
                    return f"http://{host_ip}:50021"
    except:
        pass
    
    # フォールバック: localhost
    return "http://localhost:50021"

def generate_voice_with_voicevox(text, speaker_id=10, output_path=None):
    """VOICEVOXを使用して音声を生成（雨晴はう: speaker_id=10）"""
    import requests
    import json
    import tempfile
    
    if output_path is None:
        output_path = tempfile.mktemp(suffix='.wav')
    
    # 環境に適したVOICEVOX URLを取得
    base_url = get_voicevox_url()
    
    try:
        # VOICEVOXエンジンの起動確認
        response = requests.get(f"{base_url}/speakers", timeout=5)
        if response.status_code != 200:
            raise Exception(f"VOICEVOXエンジンが起動していません (接続先: {base_url})")
        
        # 音響特徴量の生成
        query_response = requests.post(
            f"{base_url}/audio_query?text={text}&speaker={speaker_id}",
            timeout=10
        )
        query_response.raise_for_status()
        query_data = query_response.json()
        
        # 音声合成
        synthesis_response = requests.post(
            f"{base_url}/synthesis?speaker={speaker_id}",
            headers={"Content-Type": "application/json"},
            data=json.dumps(query_data),
            timeout=30
        )
        synthesis_response.raise_for_status()
        
        # 音声ファイルを保存
        with open(output_path, 'wb') as f:
            f.write(synthesis_response.content)
        
        return output_path
        
    except requests.exceptions.RequestException as e:
        raise Exception(f"VOICEVOXとの通信に失敗しました (接続先: {base_url}): {str(e)}")
    except Exception as e:
        raise Exception(f"音声生成に失敗しました: {str(e)}")

def add_multiple_voices_to_video(video_path, output_path, voices, original_volume=1.0):
    """動画に複数の音声を追加（FFmpeg直接実行版）"""
    import subprocess
    import tempfile
    
    print(f"DEBUG: FFmpeg直接実行版で音声追加開始")
    
    temp_voice_files = []
    try:
        # VOICEVOX音声を生成
        voice_files = []
        for voice in voices:
            try:
                voice_path = generate_voice_with_voicevox(voice['text'])
                temp_voice_files.append(voice_path)
                voice_files.append({
                    'path': voice_path,
                    'start_time': voice['start_time'],
                    'volume': voice['volume']
                })
            except Exception as e:
                st.warning(f"⚠️ 音声「{voice['text'][:20]}...」の生成をスキップしました: {str(e)}")
                continue
        
        if not voice_files:
            # 音声追加がない場合は元動画をそのままコピー
            import shutil
            shutil.copy2(video_path, output_path)
            return output_path
        
        # FFmpegコマンドを構築（動画ストリームはコピー、音声のみ処理）
        ffmpeg_cmd = ['ffmpeg', '-i', video_path, '-y']
        
        # 各音声ファイルを入力として追加
        for voice_file in voice_files:
            ffmpeg_cmd.extend(['-i', voice_file['path']])
        
        # フィルター構築（シンプルに音声をミックス）
        if len(voice_files) == 1:
            # 1つの音声のみ
            voice = voice_files[0]
            delay_ms = int(voice['start_time'] * 1000)  # ミリ秒に変換
            if delay_ms > 0:
                audio_filter = f'[1:a]volume={voice["volume"]},adelay={delay_ms}[voice];[0:a][voice]amix=inputs=2:duration=first[audio]'
            else:
                audio_filter = f'[1:a]volume={voice["volume"]}[voice];[0:a][voice]amix=inputs=2:duration=first[audio]'
        else:
            # 複数音声をミックス
            voice_filters = []
            for i, voice in enumerate(voice_files):
                delay_ms = int(voice['start_time'] * 1000)  # ミリ秒に変換
                if delay_ms > 0:
                    voice_filters.append(f'[{i+1}:a]volume={voice["volume"]},adelay={delay_ms}[voice{i}]')
                else:
                    voice_filters.append(f'[{i+1}:a]volume={voice["volume"]}[voice{i}]')
            
            # 全音声をミックス
            voice_labels = ''.join(f'[voice{i}]' for i in range(len(voice_files)))
            audio_filter = ';'.join(voice_filters) + f';[0:a]{voice_labels}amix=inputs={len(voice_files)+1}:duration=first[audio]'
        
        # 動画ストリームはコピー、音声のみ処理
        ffmpeg_cmd.extend([
            '-filter_complex', audio_filter,
            '-map', '0:v',  # 動画ストリームはそのままコピー
            '-map', '[audio]',  # 処理された音声
            '-c:v', 'copy',  # 動画は再エンコードしない（重要！）
            '-c:a', 'aac',
            output_path
        ])
        
        print(f"DEBUG: FFmpeg実行: {' '.join(ffmpeg_cmd)}")
        result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
        
        if result.returncode != 0:
            print(f"DEBUG: FFmpeg エラー: {result.stderr}")
            raise Exception(f"FFmpeg処理に失敗しました: {result.stderr}")
        
        print(f"DEBUG: FFmpeg成功")
        
    finally:
        # 一時ファイル削除
        for temp_file in temp_voice_files:
            try:
                os.unlink(temp_file)
            except:
                pass
    
    return output_path


def add_bgm_to_video(video_path, output_path, bgm_path=None, bgm_volume=0.5, original_volume=1.0, loop_bgm=True, bgm_start_time=0.0):
    """動画にBGMを追加（FFmpegを使用してより正確に）"""
    import subprocess
    import tempfile
    
    print(f"DEBUG BGM: FFmpeg方式でBGM追加開始")
    
    # 元の動画の情報を取得
    clip = VideoFileClip(video_path)
    original_video_duration = clip.duration
    original_fps = clip.fps
    print(f"DEBUG BGM: 元の動画 - 長さ: {original_video_duration}秒, FPS: {original_fps}")
    clip.close()
    
    if bgm_path and os.path.exists(bgm_path):
        try:
            # FFmpegでBGMを追加
            bgm_info = AudioFileClip(bgm_path)
            bgm_duration = bgm_info.duration
            bgm_info.close()
            
            ffmpeg_cmd = ['ffmpeg', '-i', video_path, '-i', bgm_path, '-y']
            
            # フィルター構築
            filter_parts = []
            
            # 動画トラック
            filter_parts.append('[0:v]copy[video]')
            
            # BGM処理
            if loop_bgm and bgm_duration < original_video_duration:
                # ループが必要な場合
                loops_needed = int(original_video_duration / bgm_duration) + 1
                filter_parts.append(f'[1:a]stream_loop={loops_needed},atrim=0:{original_video_duration},volume={bgm_volume}[bgm]')
            else:
                # ループ不要またはBGMが十分長い場合
                filter_parts.append(f'[1:a]atrim=0:{original_video_duration},volume={bgm_volume}[bgm]')
            
            # BGMの開始時間調整
            if bgm_start_time > 0.0:
                delay_ms = int(bgm_start_time * 1000)
                filter_parts[-1] = f'[1:a]atrim=0:{original_video_duration},adelay={delay_ms}|{delay_ms},volume={bgm_volume}[bgm]'
            
            # 元の音声がある場合はミックス
            if VideoFileClip(video_path).audio is not None:
                filter_parts.append(f'[0:a]volume={original_volume}[orig]')
                filter_parts.append('[orig][bgm]amix=inputs=2:duration=first[audio]')
            else:
                filter_parts.append('[bgm]acopy[audio]')
            
            # フィルターグラフを完成
            filter_complex = ';'.join(filter_parts)
            ffmpeg_cmd.extend([
                '-filter_complex', filter_complex,
                '-map', '[video]',
                '-map', '[audio]',
                '-t', str(original_video_duration),
                '-r', str(original_fps),
                '-c:v', 'libx264',
                '-c:a', 'aac',
                '-crf', '18',
                '-preset', 'slow',
                output_path
            ])
            
            print(f"DEBUG BGM: FFmpeg実行: {' '.join(ffmpeg_cmd)}")
            result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
            
            if result.returncode == 0:
                print(f"DEBUG BGM: FFmpeg成功")
                return output_path
            else:
                print(f"DEBUG BGM: FFmpeg エラー: {result.stderr}")
                # フォールバック処理へ続行
                
        except Exception as e:
            print(f"DEBUG BGM: FFmpeg方式失敗: {str(e)}")
            # フォールバック処理へ続行
        
        # フォールバック：MoviePy方式
        print("DEBUG BGM: MoviePyフォールバック方式を使用")
        clip = VideoFileClip(video_path)
        bgm = AudioFileClip(bgm_path)
        
        # BGMの音量を調整
        bgm = bgm.with_fps(44100)
        if bgm_volume != 1.0:
            bgm = bgm.with_volume_scaled(bgm_volume)
        
        # BGMをループ再生するかどうか
        if loop_bgm and bgm.duration < original_video_duration:
            # BGMをループして動画の長さに合わせる
            loops_needed = int(clip.duration / bgm.duration) + 1
            try:
                # ループ処理を行い、動画の長さに合わせてカット
                bgm = bgm.audio_loop(n=loops_needed).subclipped(0, clip.duration)
            except Exception as e:
                # ループ処理に失敗した場合は、単純に繰り返し再生
                bgm_list = []
                current_duration = 0
                while current_duration < clip.duration:
                    remaining_duration = clip.duration - current_duration
                    if remaining_duration >= bgm.duration:
                        bgm_list.append(bgm)
                        current_duration += bgm.duration
                    else:
                        bgm_list.append(bgm.subclipped(0, remaining_duration))
                        current_duration += remaining_duration
                
                if bgm_list:
                    bgm = concatenate_audioclips(bgm_list)
                else:
                    bgm = bgm.subclipped(0, min(bgm.duration, clip.duration))
        else:
            # BGMを動画の長さに合わせてカット
            bgm = bgm.subclipped(0, min(bgm.duration, clip.duration))
        
        # BGMの開始時間を適用
        if bgm_start_time > 0.0 and bgm_start_time < clip.duration:
            # 無音の音声を作成してBGMの前に追加
            silence_duration = min(bgm_start_time, clip.duration)
            remaining_duration = clip.duration - silence_duration
            
            if remaining_duration > 0:
                # BGMを残り時間に合わせてカット
                bgm = bgm.subclipped(0, min(bgm.duration, remaining_duration))
                # 無音とBGMを結合
                bgm = bgm.with_start(bgm_start_time)
            else:
                # 開始時間が動画の長さを超えている場合は無音
                bgm = None
        
        # 元の音声があるかチェック
        if clip.audio is not None:
            # 元の音声の音量を調整
            original_audio = clip.audio
            if original_volume != 1.0:
                original_audio = original_audio.with_volume_scaled(original_volume)
            # BGMと元の音声を合成
            if bgm is not None:
                final_audio = CompositeAudioClip([original_audio, bgm])
            else:
                final_audio = original_audio
        else:
            # 元の音声がない場合はBGMのみ
            if bgm is not None:
                final_audio = bgm
            else:
                final_audio = None
        
        # 動画に音声を設定（フレームレートを維持）
        if final_audio is not None:
            # 音声の長さを動画の長さに合わせる
            if final_audio.duration > original_video_duration:
                final_audio = final_audio.subclipped(0, original_video_duration)
            elif final_audio.duration < original_video_duration:
                final_audio = final_audio.with_duration(original_video_duration)
            
            # 重要：元の動画のFPSを保持
            final_clip = clip.with_audio(final_audio).with_fps(original_fps)
        else:
            final_clip = clip.with_fps(original_fps)
        
        # 動画の長さを確実に設定
        final_clip = final_clip.with_duration(original_video_duration)
        
        # 出力（元のFPSを明示的に指定）
        final_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            fps=original_fps,  # 重要：元のFPSを明示的に指定
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        
        # リソースをクリーンアップ
        clip.close()
        bgm.close()
        final_clip.close()
        if clip.audio is not None:
            original_audio.close()
        final_audio.close()
    else:
        # BGMがない場合は元の動画をそのままコピー（FPSを保持）
        clip = VideoFileClip(video_path)
        final_clip = clip.with_fps(original_fps).with_duration(original_video_duration)
        final_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            fps=original_fps,
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        clip.close()
        final_clip.close()
    
    return output_path

def add_text_to_video(video_path, output_path, telops, font_size=60):
    """動画に時間ベースのテキストオーバーレイを追加"""
    import cv2
    import numpy as np
    from PIL import Image, ImageDraw, ImageFont
    
    # OpenCVとPillowを使用してテキストを追加
    clip = VideoFileClip(video_path)
    
    def add_text_frame(get_frame, t):
        frame = get_frame(t)
        # numpy配列をPIL Imageに変換
        img = Image.fromarray(frame.astype('uint8'))
        draw = ImageDraw.Draw(img)
        
        # 日本語対応フォントを優先的に使用
        japanese_fonts = [
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "fonts/NotoSansJP-Regular.ttf",
            "NotoSansCJK-Regular.ttc"
        ]
        
        font = None
        for font_path in japanese_fonts:
            try:
                font = ImageFont.truetype(font_path, font_size)
                break
            except:
                continue
        
        if font is None:
            try:
                font = ImageFont.load_default(size=font_size)
            except:
                font = ImageFont.load_default()
        
        # 現在の時間に表示すべきテロップを描画
        for telop in telops:
            # 時間範囲内かチェック
            if telop['start_time'] <= t <= telop['end_time']:
                text = telop['text']
                position = telop['position']
                
                if not text:  # 空のテキストはスキップ
                    continue
                    
                # テキストサイズを取得
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
                
                # 位置を計算（中央寄りに調整）
                if position == "top":
                    x = (img.width - text_width) // 2
                    y = img.height // 4 - text_height // 2  # 中央寄りの上
                elif position == "bottom":
                    x = (img.width - text_width) // 2
                    y = img.height * 3 // 4 - text_height // 2  # 中央寄りの下
                else:  # center
                    x = (img.width - text_width) // 2
                    y = (img.height - text_height) // 2
                
                # テロップの色設定（デフォルトは白）
                text_color = telop.get('color', (255, 255, 255))
                
                # テキストを描画（影付き）
                draw.text((x+2, y+2), text, font=font, fill=(0, 0, 0))  # 影
                draw.text((x, y), text, font=font, fill=text_color)  # テキスト
        
        return np.array(img)
    
    # テキスト付きの動画を作成
    final_video = clip.transform(add_text_frame)
    
    # 出力（高画質設定）
    final_video.write_videofile(
        output_path, 
        codec='libx264', 
        audio_codec='aac',
        bitrate='8000k',  # 高ビットレート
        ffmpeg_params=['-crf', '18', '-preset', 'slow']  # 高画質・低圧縮
    )
    clip.close()
    final_video.close()
    
    return output_path

def extract_slides_and_notes(pptx_file):
    """PowerPointファイルからスライドと speaker notes を抽出"""
    presentation = Presentation(pptx_file)
    slides_data = []
    
    for i, slide in enumerate(presentation.slides):
        # スライドを画像として保存
        slide_image_path = tempfile.mktemp(suffix=f'_slide_{i}.png')
        
        # スライドの画像を取得するため、まずPILで空の画像を作成
        # 注意: python-pptxはスライドの直接的な画像変換をサポートしていないため、
        # ここではスライドのテキスト内容とノートのみを抽出します
        
        # スライドのテキスト内容を取得
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        
        # スピーカーノートを取得
        notes_slide = slide.notes_slide
        notes_text = ""
        if notes_slide:
            notes_text_frame = notes_slide.notes_text_frame
            if notes_text_frame:
                notes_text = notes_text_frame.text
        
        slides_data.append({
            'slide_number': i + 1,
            'slide_text': slide_text.strip(),
            'notes_text': notes_text.strip(),
            'slide_image_path': None  # 実際の画像抽出は別の方法で実装
        })
    
    return slides_data

def create_slide_images_from_pptx(pptx_file):
    """PowerPointスライドを画像ファイルに変換する（LibreOfficeを使用）"""
    import subprocess
    import shutil
    
    # 一時ディレクトリを作成
    temp_dir = tempfile.mkdtemp()
    pptx_path = os.path.join(temp_dir, "presentation.pptx")
    
    # アップロードファイルを保存
    with open(pptx_path, 'wb') as f:
        pptx_file.seek(0)
        f.write(pptx_file.read())
    
    try:
        # LibreOfficeでPDFに変換
        cmd = [
            'libreoffice', 
            '--headless', 
            '--convert-to', 'pdf',
            '--outdir', temp_dir,
            pptx_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        if result.returncode != 0:
            raise Exception(f"LibreOffice変換エラー: {result.stderr}")
        
        # PDFをPNG画像に変換（pdftoppmを使用 - より安定）
        pdf_path = os.path.join(temp_dir, "presentation.pdf")
        if not os.path.exists(pdf_path):
            raise Exception("PDF変換に失敗しました")
        
        # pdftoppmでPDFの各ページをPNGに変換
        cmd = [
            'pdftoppm',
            '-png',
            '-r', '150',  # 解像度150dpi
            pdf_path,
            os.path.join(temp_dir, 'slide')
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        if result.returncode != 0:
            # pdftoppmが失敗した場合はImageMagickを試行
            cmd = [
                'convert',
                '-density', '150',
                '-background', 'white',
                '-alpha', 'remove',
                '-quality', '90',
                pdf_path,
                os.path.join(temp_dir, 'slide-%03d.png')
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                raise Exception(f"画像変換エラー: {result.stderr}")
        
        # 生成された画像ファイルパスを取得
        slide_images = []
        
        # pdftoppmの出力形式に対応
        for file in sorted(os.listdir(temp_dir)):
            if file.startswith('slide') and file.endswith('.png'):
                image_path = os.path.join(temp_dir, file)
                if os.path.exists(image_path):
                    # 永続的な場所にコピー
                    slide_num = len(slide_images)
                    permanent_path = tempfile.mktemp(suffix=f'_slide_{slide_num}.png')
                    shutil.copy(image_path, permanent_path)
                    slide_images.append(permanent_path)
        
        if not slide_images:
            raise Exception("スライド画像の生成に失敗しました")
        
        return slide_images, temp_dir
        
    except subprocess.TimeoutExpired:
        raise Exception("変換処理がタイムアウトしました")
    except Exception as e:
        # 一時ディレクトリをクリーンアップ
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise e

def create_slide_video_with_narration(slide_image_path, narration_audio_path, duration, output_path):
    """スライド画像とナレーション音声から動画を作成"""
    from moviepy import ImageClip
    
    # スライド画像から動画クリップを作成
    slide_clip = ImageClip(slide_image_path, duration=duration)
    
    # 音声を読み込み
    audio = AudioFileClip(narration_audio_path)
    
    # 画像クリップに音声を追加
    final_clip = slide_clip.with_audio(audio)
    
    # 出力
    final_clip.write_videofile(
        output_path,
        codec='libx264',
        audio_codec='aac',
        fps=1,  # スライドなので低いFPSで十分
        ffmpeg_params=['-crf', '18', '-preset', 'fast']
    )
    
    # クリーンアップ
    slide_clip.close()
    audio.close()
    final_clip.close()
    
    return output_path

def create_text_slide_image(text, title, width=1920, height=1080):
    """テキストからスライド画像を生成"""
    # 空の画像を作成
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)
    
    try:
        # 日本語フォントを読み込み
        title_font = ImageFont.truetype("NotoSansCJK-Regular.ttc", 72)
        text_font = ImageFont.truetype("NotoSansCJK-Regular.ttc", 48)
    except:
        try:
            title_font = ImageFont.load_default(size=72)
            text_font = ImageFont.load_default(size=48)
        except:
            title_font = ImageFont.load_default()
            text_font = ImageFont.load_default()
    
    # タイトルを描画
    title_bbox = draw.textbbox((0, 0), title, font=title_font)
    title_width = title_bbox[2] - title_bbox[0]
    title_x = (width - title_width) // 2
    draw.text((title_x, 100), title, fill='black', font=title_font)
    
    # テキストを描画（改行対応）
    lines = text.split('\n')
    y_offset = 250
    line_height = 60
    
    for line in lines:
        if not line.strip():
            continue
        
        # 長い行を自動改行
        words = line.split()
        current_line = ""
        
        for word in words:
            test_line = current_line + word + " "
            bbox = draw.textbbox((0, 0), test_line, font=text_font)
            test_width = bbox[2] - bbox[0]
            
            if test_width > width - 200:  # マージンを考慮
                if current_line:
                    # 現在の行を描画
                    text_bbox = draw.textbbox((0, 0), current_line, font=text_font)
                    text_width = text_bbox[2] - text_bbox[0]
                    text_x = (width - text_width) // 2
                    draw.text((text_x, y_offset), current_line, fill='black', font=text_font)
                    y_offset += line_height
                current_line = word + " "
            else:
                current_line = test_line
        
        # 残りのテキストを描画
        if current_line.strip():
            text_bbox = draw.textbbox((0, 0), current_line, font=text_font)
            text_width = text_bbox[2] - text_bbox[0]
            text_x = (width - text_width) // 2
            draw.text((text_x, y_offset), current_line, fill='black', font=text_font)
            y_offset += line_height
    
    # 画像を保存
    output_path = tempfile.mktemp(suffix='.png')
    img.save(output_path)
    
    return output_path

def combine_videos(video_paths, output_path):
    """複数の動画を結合する"""
    clips = []
    try:
        for video_path in video_paths:
            # ファイルの存在確認
            if not os.path.exists(video_path):
                raise FileNotFoundError(f"ファイルが見つかりません: {video_path}")
            
            clip = VideoFileClip(video_path)
            clips.append(clip)
        
        # 動画を結合
        final_clip = concatenate_videoclips(clips, method="compose")
        
        # 出力（高画質設定）
        final_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        
        return output_path
        
    finally:
        # リソースをクリーンアップ
        for clip in clips:
            try:
                clip.close()
            except:
                pass
        try:
            final_clip.close()
        except:
            pass

# メインインターface
if tool == "ショート動画変換":
    uploaded_file = st.file_uploader(
        "動画ファイルをアップロードしてください",
        type=['mp4', 'avi', 'mov', 'mkv'],
        help="MP4, AVI, MOV, MKV形式の動画ファイルをサポートしています"
    )
elif tool == "動画結合":
    uploaded_files = st.file_uploader(
        "結合する動画ファイルを複数選択してください",
        type=['mp4', 'avi', 'mov', 'mkv'],
        accept_multiple_files=True,
        help="MP4, AVI, MOV, MKV形式の動画ファイルをサポートしています"
    )
else:  # パワポナレーション動画
    uploaded_pptx = st.file_uploader(
        "PowerPointファイルをアップロードしてください",
        type=['pptx', 'ppt'],
        help="PowerPoint形式のファイル（.pptx, .ppt）をサポートしています"
    )

if tool == "ショート動画変換" and uploaded_file is not None:
    # 一時ファイルに保存
    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp_file:
        tmp_file.write(uploaded_file.read())
        input_video_path = tmp_file.name
    
    # 動画情報を表示
    try:
        clip = VideoFileClip(input_video_path)
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("解像度", f"{clip.w}x{clip.h}")
        with col2:
            st.metric("時間", f"{clip.duration:.1f}秒")
        with col3:
            st.metric("FPS", f"{clip.fps:.1f}")
        
        clip.close()
        
        # オプション設定
        st.subheader("設定オプション")
        
        # 動画トリミング設定
        st.subheader("✂️ 動画トリミング設定")
        trim_video = st.checkbox("不要な部分を削除する")
        
        if trim_video:
            col1, col2 = st.columns(2)
            with col1:
                start_time = st.number_input(
                    "開始時間（秒）", 
                    min_value=0, 
                    max_value=int(clip.duration), 
                    value=0, 
                    step=1,
                    help="この時間から動画を開始します"
                )
            with col2:
                end_time = st.number_input(
                    "終了時間（秒）", 
                    min_value=start_time + 1, 
                    max_value=int(clip.duration), 
                    value=int(clip.duration), 
                    step=1,
                    help="この時間で動画を終了します"
                )
            
            if end_time > start_time:
                trimmed_duration = end_time - start_time
                st.info(f"⏱️ トリミング後の長さ: {trimmed_duration:.1f}秒")
            else:
                st.error("終了時間は開始時間より後に設定してください")
        
        st.subheader("📏 拡大倍率設定")
        keep_original_size = st.checkbox("元の動画サイズを維持する", help="チェックすると、動画は9:16にリサイズされず、元の解像度のまま処理されます。")
        
        if not keep_original_size:
            scale_factor = st.slider(
                "動画の拡大倍率",
                min_value=0.5,
                max_value=5.0,
                value=1.0,
                step=0.1,
                help="1.0 = 原寸大、3.0 = 300%拡大。大きくするほどズームインされます。"
            )
            
            # 倍率の説明表示
            if scale_factor < 1.0:
                st.info(f"📉 縮小表示: {scale_factor*100:.0f}% (より多くの映像が見えます)")
            elif scale_factor == 1.0:
                st.info("📊 原寸大表示: 100% (元の動画のまま)")
            else:
                st.info(f"📈 拡大表示: {scale_factor*100:.0f}% (ズームイン効果)")
        else:
            scale_factor = 1.0 # 元のサイズを維持する場合はスケールは1.0固定
        
        # テキストオーバーレイ設定
        add_text = st.checkbox("テキストを追加する")
        
        if add_text:
            st.subheader("テロップ設定")
            
            # セッション状態でテロップリストを管理
            if 'telops' not in st.session_state:
                st.session_state.telops = []
            
            font_size = st.slider("フォントサイズ", 30, 120, 60)
            
            # フォント色の選択
            st.subheader("カラー設定")
            color_options = {
                "白": (255, 255, 255),
                "黒": (0, 0, 0),
                "赤": (255, 0, 0),
                "青": (0, 0, 255),
                "緑": (0, 255, 0),
                "黄": (255, 255, 0),
                "マゼンタ": (255, 0, 255),
                "シアン": (0, 255, 255),
                "オレンジ": (255, 165, 0)
            }
            selected_color_name = st.selectbox("フォント色", list(color_options.keys()), index=0)
            selected_color = color_options[selected_color_name]
            
            # 新しいテロップを追加するフォーム
            with st.expander("新しいテロップを追加", expanded=True):
                col1, col2, col3, col4, col5, col6 = st.columns([2, 1, 1, 1, 1, 1])
                
                with col1:
                    new_text = st.text_input("テロップテキスト", key="new_telop_text")
                with col2:
                    new_position = st.selectbox("位置", ["top", "center", "bottom"], key="new_telop_position")
                with col3:
                    new_start = st.number_input("開始(秒)", min_value=0, value=0, step=1, key="new_telop_start")
                with col4:
                    new_end = st.number_input("終了(秒)", min_value=1, value=5, step=1, key="new_telop_end")
                with col5:
                    new_color_name = st.selectbox("色", list(color_options.keys()), index=0, key="new_telop_color")
                with col6:
                    if st.button("追加", type="primary"):
                        if new_text.strip():
                            new_telop = {
                                "text": new_text,
                                "position": new_position,
                                "start_time": new_start,
                                "end_time": new_end,
                                "color": color_options[new_color_name]
                            }
                            st.session_state.telops.append(new_telop)
                            st.rerun()
            
            # 既存のテロップを表示・編集
            if st.session_state.telops:
                st.subheader("登録済みテロップ")
                
                for i, telop in enumerate(st.session_state.telops):
                    with st.container():
                        col1, col2, col3, col4, col5, col6, col7 = st.columns([2, 1, 1, 1, 1, 1, 0.5])
                        
                        with col1:
                            st.text_input(f"テキスト {i+1}", value=telop["text"], key=f"telop_text_{i}", disabled=True)
                        with col2:
                            st.text(telop["position"])
                        with col3:
                            st.text(f"{telop['start_time']}秒")
                        with col4:
                            st.text(f"{telop['end_time']}秒")
                        with col5:
                            duration = telop['end_time'] - telop['start_time']
                            st.text(f"{duration:.1f}秒間")
                        with col6:
                            # 色を表示
                            color = telop.get('color', (255, 255, 255))
                            color_name = next((name for name, rgb in color_options.items() if rgb == color), "白")
                            st.text(color_name)
                        with col7:
                            if st.button("🗑️", key=f"delete_telop_{i}", help="削除"):
                                st.session_state.telops.pop(i)
                                st.rerun()
                
                if st.button("全てのテロップをクリア", type="secondary"):
                    st.session_state.telops = []
                    st.rerun()
        
        # 音声合成設定（雨晴はう）
        st.subheader("🎤 音声合成設定（雨晴はう）")
        
        # VOICEVOXの設定情報を表示
        st.info("""
        💡 **VOICEVOX設定について (WSL環境)**:
        1. Windows側でVOICEVOXを起動してください
        2. VOICEVOXの設定で「外部からのアクセスを許可」を有効にしてください
        3. Windowsファイアウォールで port 50021 を開放してください
        """)
        
        add_voice = st.checkbox("雨晴はうの音声を追加する")
        
        if add_voice:
            # セッション状態で音声リストを管理
            if 'voices' not in st.session_state:
                st.session_state.voices = []
            
            # 新しい音声を追加するフォーム
            with st.expander("新しい音声を追加", expanded=True):
                col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                
                with col1:
                    new_voice_text = st.text_area(
                        "読み上げテキスト",
                        placeholder="雨晴はうに読み上げてもらいたいテキストを入力してください",
                        key="new_voice_text",
                        height=100
                    )
                with col2:
                    new_voice_start = st.number_input("開始時間(秒)", min_value=0, value=0, step=1, key="new_voice_start")
                with col3:
                    new_voice_volume = st.slider("音量", 0.0, 1.0, 0.8, 0.1, key="new_voice_volume")
                with col4:
                    st.write("") # スペース調整
                    if st.button("🔊 プレビュー", key="preview_voice"):
                        if new_voice_text.strip():
                            try:
                                with st.spinner("音声を生成中..."):
                                    voice_path = generate_voice_with_voicevox(new_voice_text)
                                    st.audio(voice_path)
                                    os.unlink(voice_path)
                                    st.success("✅ 音声生成成功！")
                            except Exception as e:
                                st.error(f"❌ 音声生成エラー: {str(e)}")
                    
                    if st.button("追加", type="primary", key="add_voice"):
                        if new_voice_text.strip():
                            new_voice = {
                                "text": new_voice_text,
                                "start_time": new_voice_start,
                                "volume": new_voice_volume
                            }
                            st.session_state.voices.append(new_voice)
                            st.rerun()
            
            # 既存の音声を表示・編集
            if st.session_state.voices:
                st.subheader("登録済み音声")
                
                for i, voice in enumerate(st.session_state.voices):
                    with st.container():
                        col1, col2, col3, col4 = st.columns([3, 1, 1, 0.5])
                        
                        with col1:
                            # テキストを1行で表示（長い場合は省略）
                            display_text = voice["text"][:50] + "..." if len(voice["text"]) > 50 else voice["text"]
                            st.text_input(f"音声 {i+1}", value=display_text, key=f"voice_text_{i}", disabled=True)
                        with col2:
                            st.text(f"{voice['start_time']}秒から")
                        with col3:
                            st.text(f"音量: {voice['volume']}")
                        with col4:
                            if st.button("🗑️", key=f"delete_voice_{i}", help="削除"):
                                st.session_state.voices.pop(i)
                                st.rerun()
                
                if st.button("全ての音声をクリア", type="secondary"):
                    st.session_state.voices = []
                    st.rerun()
        
        # BGM設定
        st.subheader("🎵 BGM設定")
        add_bgm = st.checkbox("BGMを追加する")
        
        if add_bgm:
            bgm_file = st.file_uploader(
                "BGMファイルをアップロード",
                type=['mp3', 'wav', 'aac', 'm4a', 'ogg'],
                help="MP3, WAV, AAC, M4A, OGG形式の音声ファイルをサポートしています"
            )
            
            col1, col2 = st.columns(2)
            with col1:
                bgm_volume = st.slider("BGM音量", 0.0, 1.0, 0.3, 0.1)
            with col2:
                original_volume = st.slider("元音声音量", 0.0, 1.0, 0.7, 0.1)
            
            loop_bgm = st.checkbox("BGMをループ再生", value=True)
            
            if bgm_file:
                st.audio(bgm_file)
        
        # 変換ボタン
        if st.button("ショート動画に変換", type="primary"):
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                # BGMファイルを一時保存
                bgm_path = None
                if add_bgm and bgm_file is not None:
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as tmp_bgm:
                        tmp_bgm.write(bgm_file.read())
                        bgm_path = tmp_bgm.name
                
                # Step 1: 動画をショート形式にリサイズ
                status_text.text("動画をリサイズ中...")
                progress_bar.progress(20)
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='_resized.mp4') as tmp_resized:
                    resized_video_path = tmp_resized.name
                
                # トリミングパラメータを設定
                trim_start = start_time if trim_video else None
                trim_end = end_time if trim_video else None
                
                resize_video_to_shorts(input_video_path, resized_video_path, scale_factor, trim_start, trim_end, keep_original_size)
                progress_bar.progress(40)
                
                # Step 2: テキストを追加（オプション）
                current_video_path = resized_video_path
                if add_text and st.session_state.telops:
                    status_text.text("テキストを追加中...")
                    progress_bar.progress(60)
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix='_with_text.mp4') as tmp_text:
                        text_video_path = tmp_text.name
                    
                    add_text_to_video(current_video_path, text_video_path, st.session_state.telops, font_size)
                    os.unlink(current_video_path)
                    current_video_path = text_video_path
                
                # Step 3: 音声合成を追加（オプション）
                if add_voice and st.session_state.voices:
                    status_text.text("雨晴はうの音声を生成・追加中...")
                    progress_bar.progress(60)
                    
                    try:
                        # 複数音声を動画に追加
                        with tempfile.NamedTemporaryFile(delete=False, suffix='_with_voices.mp4') as tmp_voice:
                            voice_video_path = tmp_voice.name
                        
                        add_multiple_voices_to_video(
                            current_video_path,
                            voice_video_path,
                            st.session_state.voices,
                            1.0  # 元音声音量
                        )
                        
                        os.unlink(current_video_path)
                        current_video_path = voice_video_path
                        
                    except Exception as e:
                        st.warning(f"⚠️ 音声合成をスキップしました: {str(e)}")
                
                # Step 4: BGMを追加（オプション）
                if add_bgm and bgm_path:
                    status_text.text("BGMを追加中...")
                    progress_bar.progress(80)
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix='_final.mp4') as tmp_final:
                        final_video_path = tmp_final.name
                    
                    add_bgm_to_video(
                        current_video_path, 
                        final_video_path, 
                        bgm_path, 
                        bgm_volume, 
                        original_volume, 
                        loop_bgm
                    )
                    os.unlink(current_video_path)
                    os.unlink(bgm_path)
                else:
                    final_video_path = current_video_path
                
                progress_bar.progress(100)
                status_text.text("変換完了！")
                
                # プレビュー表示
                st.subheader("📹 プレビュー")
                st.video(final_video_path, width=400)
                
                # ダウンロードボタンを表示
                with open(final_video_path, 'rb') as file:
                    st.download_button(
                        label="📱 ショート動画をダウンロード",
                        data=file.read(),
                        file_name=f"shorts_{uploaded_file.name}",
                        mime="video/mp4"
                    )
                
                # 一時ファイルをクリーンアップ
                os.unlink(input_video_path)
                os.unlink(final_video_path)
                
                st.success("✅ 変換が完了しました！ダウンロードボタンをクリックして保存してください。")
                
            except Exception as e:
                st.error(f"❌ エラーが発生しました: {str(e)}")
                # エラー時のクリーンアップ
                try:
                    os.unlink(input_video_path)
                except:
                    pass
    
    except Exception as e:
        st.error(f"❌ 動画ファイルの読み込みに失敗しました: {str(e)}")
        os.unlink(input_video_path)

elif tool == "動画結合" and uploaded_files:
    if len(uploaded_files) < 2:
        st.warning("⚠️ 2つ以上の動画ファイルを選択してください。")
    else:
        st.success(f"✅ {len(uploaded_files)}個の動画ファイルが選択されました。")
        
        # 動画情報を表示
        st.subheader("📹 選択された動画")
        total_duration = 0
        
        for i, file in enumerate(uploaded_files):
            # ファイルを一時保存
            file.seek(0)  # ファイルポインタを先頭に戻す
            temp_path = tempfile.mktemp(suffix=f'_preview_{i}.mp4')
            
            try:
                with open(temp_path, 'wb') as tmp_file:
                    file_bytes = file.getvalue()
                    tmp_file.write(file_bytes)
                    tmp_file.flush()
                    os.fsync(tmp_file.fileno())
                
                # ファイルサイズを確認
                if os.path.getsize(temp_path) == 0:
                    st.error(f"❌ {file.name} のサイズが0です")
                    continue
                
            except Exception as e:
                st.error(f"❌ {file.name}の保存に失敗しました: {str(e)}")
                continue
            
            try:
                clip = VideoFileClip(temp_path)
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    st.text(f"{i+1}. {file.name}")
                with col2:
                    st.text(f"{clip.w}x{clip.h}")
                with col3:
                    st.text(f"{clip.duration:.1f}秒")
                with col4:
                    st.text(f"{clip.fps:.1f} FPS")
                
                total_duration += clip.duration
                clip.close()
                os.unlink(temp_path)
            except Exception as e:
                st.error(f"❌ {file.name}の読み込みに失敗しました: {str(e)}")
                try:
                    os.unlink(temp_path)
                except:
                    pass
        
        st.info(f"📊 結合後の総時間: {total_duration:.1f}秒")
        
        # 結合ボタン
        if st.button("動画を結合", type="primary"):
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                status_text.text("動画を結合中...")
                progress_bar.progress(20)
                
                # 一時ファイルに保存
                temp_paths = []
                for i, file in enumerate(uploaded_files):
                    file.seek(0)  # ファイルポインタを先頭に戻す
                    
                    # 一意のファイル名を生成
                    temp_path = tempfile.mktemp(suffix=f'_video_{i}.mp4')
                    
                    try:
                        with open(temp_path, 'wb') as tmp_file:
                            file_bytes = file.getvalue()
                            tmp_file.write(file_bytes)
                            tmp_file.flush()
                            os.fsync(tmp_file.fileno())
                        
                        # ファイルサイズを確認
                        if os.path.getsize(temp_path) == 0:
                            raise ValueError(f"ファイル {file.name} のサイズが0です")
                        
                        temp_paths.append(temp_path)
                        
                    except Exception as e:
                        st.error(f"❌ {file.name}の保存に失敗しました: {str(e)}")
                        # 失敗したファイルがあればクリーンアップして終了
                        for path in temp_paths:
                            try:
                                os.unlink(path)
                            except:
                                pass
                        raise
                
                progress_bar.progress(50)
                
                # 動画を結合
                with tempfile.NamedTemporaryFile(delete=False, suffix='_combined.mp4') as tmp_output:
                    output_path = tmp_output.name
                
                combine_videos(temp_paths, output_path)
                progress_bar.progress(100)
                status_text.text("結合完了！")
                
                # プレビュー表示
                st.subheader("📹 結合された動画")
                st.video(output_path, width=400)
                
                # ダウンロードボタンを表示
                with open(output_path, 'rb') as file:
                    st.download_button(
                        label="📱 結合動画をダウンロード",
                        data=file.read(),
                        file_name=f"combined_{len(uploaded_files)}_videos.mp4",
                        mime="video/mp4"
                    )
                
                # 一時ファイルをクリーンアップ
                for temp_path in temp_paths:
                    os.unlink(temp_path)
                os.unlink(output_path)
                
                st.success("✅ 結合が完了しました！ダウンロードボタンをクリックして保存してください。")
                
            except Exception as e:
                st.error(f"❌ エラーが発生しました: {str(e)}")
                # エラー時のクリーンアップ
                for temp_path in temp_paths:
                    try:
                        os.unlink(temp_path)
                    except:
                        pass

elif tool == "パワポナレーション動画" and uploaded_pptx is not None:
    try:
        # PowerPointファイルからスライドとノートを抽出
        with st.spinner("PowerPointファイルを解析中..."):
            slides_data = extract_slides_and_notes(uploaded_pptx)
        
        if not slides_data:
            st.error("❌ スライドが見つかりませんでした。")
        else:
            st.success(f"✅ {len(slides_data)}枚のスライドが見つかりました。")
            
            # スライド情報を表示
            st.subheader("📋 抽出されたスライド情報")
            
            for i, slide in enumerate(slides_data):
                with st.expander(f"スライド {slide['slide_number']}: {slide['slide_text'][:50]}..."):
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        st.markdown("**スライド内容:**")
                        st.text_area("", value=slide['slide_text'], height=100, key=f"slide_text_{i}", disabled=True)
                    
                    with col2:
                        st.markdown("**ノート（読み上げ内容）:**")
                        if slide['notes_text']:
                            st.text_area("", value=slide['notes_text'], height=100, key=f"notes_text_{i}", disabled=True)
                        else:
                            st.info("ノートが設定されていません")
            
            # 設定オプション
            st.subheader("🎛️ 動画設定")
            
            col1, col2 = st.columns(2)
            
            with col1:
                slide_duration = st.slider(
                    "各スライドの表示時間（秒）",
                    min_value=3,
                    max_value=30,
                    value=10,
                    help="ノートがある場合は音声の長さに自動調整されます"
                )
            
            with col2:
                voice_speed = st.slider(
                    "読み上げ速度",
                    min_value=0.5,
                    max_value=2.0,
                    value=1.0,
                    step=0.1,
                    help="音声の読み上げ速度を調整します"
                )
            
            # 変換ボタン
            if st.button("ナレーション動画を作成", type="primary"):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                try:
                    # Step 1: PowerPointスライドを画像に変換
                    status_text.text("スライドを画像に変換中...")
                    progress_bar.progress(10)
                    
                    # PowerPointファイルから実際のスライド画像を抽出
                    try:
                        slide_image_paths, temp_conversion_dir = create_slide_images_from_pptx(uploaded_pptx)
                        use_real_slides = True
                        st.info(f"✅ {len(slide_image_paths)}枚のスライド画像を抽出しました")
                    except Exception as e:
                        st.warning(f"⚠️ スライド画像の抽出に失敗しました。テキストベースのスライドを使用します: {str(e)}")
                        slide_image_paths = []
                        use_real_slides = False
                    
                    slide_videos = []
                    temp_files = []
                    
                    for i, slide in enumerate(slides_data):
                        progress = 10 + (i / len(slides_data)) * 80
                        progress_bar.progress(int(progress))
                        status_text.text(f"スライド {i+1}/{len(slides_data)} を処理中...")
                        
                        # スライド画像を取得（実際のスライドまたはテキストベース）
                        if use_real_slides and i < len(slide_image_paths):
                            slide_image_path = slide_image_paths[i]
                        else:
                            # フォールバック: テキストベースのスライド生成
                            slide_image_path = create_text_slide_image(
                                slide['slide_text'], 
                                f"スライド {slide['slide_number']}"
                            )
                            temp_files.append(slide_image_path)
                        
                        # ナレーション音声を生成（ノートがある場合）
                        if slide['notes_text'].strip():
                            try:
                                voice_path = generate_voice_with_voicevox(slide['notes_text'])
                                temp_files.append(voice_path)
                                
                                # 音声の長さを取得
                                audio_clip = AudioFileClip(voice_path)
                                duration = max(audio_clip.duration, 3)  # 最低3秒
                                audio_clip.close()
                                
                            except Exception as e:
                                st.warning(f"⚠️ スライド{i+1}の音声生成に失敗: {str(e)}")
                                voice_path = None
                                duration = slide_duration
                        else:
                            voice_path = None
                            duration = slide_duration
                        
                        # スライド動画を作成
                        slide_video_path = tempfile.mktemp(suffix=f'_slide_video_{i}.mp4')
                        temp_files.append(slide_video_path)
                        
                        if voice_path:
                            create_slide_video_with_narration(slide_image_path, voice_path, duration, slide_video_path)
                        else:
                            # 音声なしの場合
                            from moviepy import ImageClip
                            clip = ImageClip(slide_image_path, duration=duration)
                            clip.write_videofile(
                                slide_video_path,
                                codec='libx264',
                                fps=1,
                                ffmpeg_params=['-crf', '18', '-preset', 'fast']
                            )
                            clip.close()
                        
                        slide_videos.append(slide_video_path)
                    
                    # Step 2: 全スライド動画を結合
                    status_text.text("動画を結合中...")
                    progress_bar.progress(90)
                    
                    final_output_path = tempfile.mktemp(suffix='_presentation_video.mp4')
                    temp_files.append(final_output_path)
                    
                    combine_videos(slide_videos, final_output_path)
                    
                    progress_bar.progress(100)
                    status_text.text("変換完了！")
                    
                    # プレビュー表示
                    st.subheader("📹 作成されたナレーション動画")
                    st.video(final_output_path, width=600)
                    
                    # ダウンロードボタンを表示
                    with open(final_output_path, 'rb') as file:
                        st.download_button(
                            label="📱 ナレーション動画をダウンロード",
                            data=file.read(),
                            file_name=f"presentation_{uploaded_pptx.name.split('.')[0]}.mp4",
                            mime="video/mp4"
                        )
                    
                    st.success("✅ 動画変換が完了しました！ダウンロードボタンをクリックして保存してください。")
                    
                except Exception as e:
                    st.error(f"❌ エラーが発生しました: {str(e)}")
                finally:
                    # 一時ファイルをクリーンアップ
                    for temp_file in temp_files:
                        try:
                            os.unlink(temp_file)
                        except:
                            pass
                    
                    # スライド画像の一時ファイルもクリーンアップ
                    if use_real_slides:
                        for slide_image in slide_image_paths:
                            try:
                                os.unlink(slide_image)
                            except:
                                pass
                        try:
                            import shutil
                            shutil.rmtree(temp_conversion_dir, ignore_errors=True)
                        except:
                            pass
    
    except Exception as e:
        st.error(f"❌ PowerPointファイルの解析に失敗しました: {str(e)}")

# 使用方法の説明
if tool == "ショート動画変換":
    st.sidebar.header("📖 使用方法")
    st.sidebar.markdown("""
    1. **動画をアップロード**: MP4、AVI、MOV、MKV形式の動画ファイルを選択
    2. **テキスト追加（オプション）**: テロップを追加したい場合はチェック
    3. **変換実行**: 「ショート動画に変換」ボタンをクリック
    4. **ダウンロード**: 変換完了後、ダウンロードボタンでファイルを保存
    
    **YouTubeショート仕様**:
    - 解像度: 1080x1920 (9:16)
    - 最大時間: 60秒
    """)
elif tool == "動画結合":
    st.sidebar.header("📖 使用方法")
    st.sidebar.markdown("""
    1. **動画を複数選択**: MP4、AVI、MOV、MKV形式の動画ファイルを2つ以上選択
    2. **結合実行**: 「動画を結合」ボタンをクリック
    3. **ダウンロード**: 結合完了後、ダウンロードボタンでファイルを保存
    
    **注意事項**:
    - 動画は選択した順番で結合されます
    - 異なる解像度の動画も結合可能です
    """)
else:  # パワポナレーション動画
    st.sidebar.header("📖 使用方法")
    st.sidebar.markdown("""
    1. **PowerPointファイルをアップロード**: .pptx, .ppt形式のファイルを選択
    2. **スライド確認**: 抽出されたスライド内容とノートを確認
    3. **設定調整**: スライド表示時間と読み上げ速度を調整
    4. **変換実行**: 「ナレーション動画を作成」ボタンをクリック
    5. **ダウンロード**: 完成した動画をダウンロード
    
    **重要**:
    - スピーカーノートが読み上げられます
    - ノートがないスライドは設定時間で表示されます
    """)

st.sidebar.header("ℹ️ 対応形式")
if tool == "ショート動画変換" or tool == "動画結合":
    st.sidebar.markdown("""
    **動画入力形式**:
    - MP4
    - AVI
    - MOV
    - MKV

    **BGM入力形式**:
    - MP3
    - WAV
    - AAC
    - M4A
    - OGG

    **出力形式**:
    - MP4 (H.264)
    """)
else:  # パワポナレーション動画
    st.sidebar.markdown("""
    **PowerPoint入力形式**:
    - PPTX
    - PPT

    **出力形式**:
    - MP4 (H.264)
    - 解像度: 1920x1080 (16:9)
    - 音声: AAC (VOICEVOX)
    """)