def extract_slides_from_pptx(pptx_file, output_dir):
    """PowerPointファイルを画像に変換して抽出"""
    import subprocess
    import tempfile
    import os
    from PIL import Image
    
    try:
        print("DEBUG: Starting PowerPoint to image conversion...")
        
        # 一時ファイルとディレクトリの作成
        temp_dir = tempfile.mkdtemp()
        temp_pptx_path = os.path.join(temp_dir, "presentation.pptx")
        
        # PowerPointファイルを一時保存
        with open(temp_pptx_path, 'wb') as f:
            f.write(pptx_file.getvalue() if hasattr(pptx_file, 'getvalue') else pptx_file)
        
        print(f"DEBUG: Saved PowerPoint file to: {temp_pptx_path}")
        
        # LibreOfficeを使用してPowerPointを画像に変換
        libreoffice_cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'png',
            '--outdir', temp_dir,
            temp_pptx_path
        ]
        
        print(f"DEBUG: Running LibreOffice command: {' '.join(libreoffice_cmd)}")
        
        try:
            result = subprocess.run(libreoffice_cmd, capture_output=True, text=True, timeout=60)
            print(f"DEBUG: LibreOffice return code: {result.returncode}")
            print(f"DEBUG: LibreOffice stdout: {result.stdout}")
            print(f"DEBUG: LibreOffice stderr: {result.stderr}")
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
                
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out")
        except FileNotFoundError:
            raise Exception("LibreOffice not found. Please install LibreOffice in the container.")
        
        # 変換された画像ファイルを探してリサイズ
        slide_images = []
        png_files = [f for f in os.listdir(temp_dir) if f.endswith('.png')]
        png_files.sort()
        
        print(f"DEBUG: Found {len(png_files)} PNG files: {png_files}")
        
        if not png_files:
            # LibreOfficeが失敗した場合のフォールバック
            print("DEBUG: No PNG files found, using fallback method...")
            return create_fallback_slides(pptx_file, output_dir)
        
        for i, png_file in enumerate(png_files):
            src_path = os.path.join(temp_dir, png_file)
            dst_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            
            # 画像を標準サイズ（1920x1080）にリサイズ
            with Image.open(src_path) as img:
                # 横型サイズ (16:9) にリサイズ
                img_resized = img.resize((1920, 1080), Image.Resampling.LANCZOS)
                img_resized.save(dst_path, 'PNG')
                slide_images.append(dst_path)
                print(f"DEBUG: Created slide image: {dst_path}")
        
        # 一時ファイルをクリーンアップ
        import shutil
        shutil.rmtree(temp_dir, ignore_errors=True)
        
        return slide_images, len(slide_images)
        
    except Exception as e:
        print(f"DEBUG: PowerPoint conversion error: {e}")
        # エラー時はフォールバック手法を使用
        return create_fallback_slides(pptx_file, output_dir)

def create_fallback_slides(pptx_file, output_dir):
    """フォールバック: シンプルなスライド作成"""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    
    try:
        prs = Presentation(pptx_file)
        slide_images = []
        
        for i, slide in enumerate(prs.slides):
            # シンプルなスライド画像を作成
            img = Image.new('RGB', (1920, 1080), 'white')
            draw = ImageDraw.Draw(img)
            
            # デフォルトフォントを使用
            try:
                font_title = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 80)
                font_content = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 50)
            except:
                font_title = ImageFont.load_default()
                font_content = ImageFont.load_default()
            
            # スライドタイトル
            title = f"Slide {i+1}"
            title_bbox = draw.textbbox((0, 0), title, font=font_title)
            title_x = (1920 - (title_bbox[2] - title_bbox[0])) // 2
            draw.text((title_x, 200), title, fill='black', font=font_title)
            
            # スライドコンテンツを取得
            content_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content_lines.extend(shape.text.strip().split('\n'))
            
            if not content_lines:
                content_lines = [f"Slide {i+1} Content"]
            
            # コンテンツを描画
            y_pos = 400
            for line in content_lines[:8]:  # 最大8行
                if line.strip():
                    line_bbox = draw.textbbox((0, 0), line, font=font_content)
                    line_x = (1920 - (line_bbox[2] - line_bbox[0])) // 2
                    draw.text((line_x, y_pos), line, fill='black', font=font_content)
                    y_pos += 70
            
            # 画像を保存
            slide_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            img.save(slide_path)
            slide_images.append(slide_path)
        
        return slide_images, len(prs.slides)
        
    except Exception as e:
        raise Exception(f"Fallback slide creation failed: {str(e)}")