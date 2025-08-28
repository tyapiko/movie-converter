import streamlit as st
import tempfile
import os
import time
import glob
from moviepy import VideoFileClip, AudioFileClip, CompositeAudioClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont
import numpy as np
import cv2
from pptx import Presentation
from pptx.util import Inches
import io


st.set_page_config(
    page_title="動画編集ツール",
    page_icon="🎬",
    layout="wide"
)

# サイドバーでツール選択
st.sidebar.title("🛠️ ツール選択")
tool = st.sidebar.radio(
    "使用するツール",
    ["ショート動画変換", "動画結合", "スライド動画作成"]
)

if tool == "ショート動画変換":
    st.title("🎬 ショート動画コンバーター")
    st.markdown("動画をアップロードして、YouTubeショート向けの縦型動画に変換しましょう！")
elif tool == "動画結合":
    st.title("🔗 動画結合ツール")
    st.markdown("複数のショート動画を選択して、長時間動画に結合しましょう！")
else:  # スライド動画作成
    st.title("📊 スライド動画作成ツール")
    st.markdown("PowerPointスライドと音楽から、歌詞付きの横型動画(16:9)を作成しましょう！")

def resize_video_to_shorts(video_path, output_path, scale_factor=1.0, start_time=None, end_time=None):
    """動画をYouTubeショート形式(9:16)にリサイズ（全体表示、パディング付き）"""
    import subprocess
    
    # 元の動画情報を取得
    clip = VideoFileClip(video_path)
    original_width, original_height = clip.size
    original_ratio = original_width / original_height
    clip.close()
    
    # YouTubeショートの推奨解像度: 1080x1920 (9:16)
    target_width = 1080
    target_height = 1920
    target_ratio = target_width / target_height
    
    # 基本スケール計算（ターゲット枠に収まるサイズ）
    if original_ratio > target_ratio:
        # 横長の場合、幅をターゲット幅に合わせる
        base_width = target_width
        base_height = int(target_width / original_ratio)
    else:
        # 縦長の場合、高さをターゲット高さに合わせる
        base_height = target_height
        base_width = int(target_height * original_ratio)
    
    # scale_factorを適用（拡大倍率による調整）
    final_width = int(base_width * scale_factor)
    final_height = int(base_height * scale_factor)
    
    # 拡大倍率が1.0より大きい場合、動画がターゲットフレームからはみ出すのは正常
    # パディングエラーを避けるため、最小サイズは1ピクセル以上を保証
    final_width = max(1, final_width)
    final_height = max(1, final_height)
    
    # FFmpegコマンドで動画変換
    ffmpeg_cmd = ['ffmpeg', '-i', video_path]
    
    # トリミングが指定されている場合
    if start_time is not None and end_time is not None:
        ffmpeg_cmd.extend(['-ss', str(start_time), '-t', str(end_time - start_time)])
    
    # ビデオフィルターを構築
    if scale_factor > 1.0:
        # 拡大時：スケール→中央クロップ→パディング
        vf = f'scale={final_width}:{final_height},crop={min(final_width, target_width)}:{min(final_height, target_height)},pad={target_width}:{target_height}:(ow-iw)/2:(oh-ih)/2:black'
    else:
        # 縮小時：スケール→パディング
        vf = f'scale={final_width}:{final_height},pad={target_width}:{target_height}:(ow-iw)/2:(oh-ih)/2:black'
    
    ffmpeg_cmd.extend([
        '-vf', vf,
        '-c:v', 'libx264',
        '-c:a', 'aac',
        '-b:v', '8000k',
        '-crf', '18',
        '-preset', 'slow',
        '-y',  # overwrite output file
        output_path
    ])
    
    try:
        subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
        return output_path
    except subprocess.CalledProcessError as e:
        raise Exception(f"FFmpeg処理でエラーが発生しました: {e.stderr.decode()}")
    except FileNotFoundError:
        raise Exception("FFmpegが見つかりません。システムにFFmpegがインストールされていることを確認してください。")

def get_voicevox_url():
    """VOICEVOX接続URLを取得（環境変数を優先）"""
    import os
    
    # 環境変数から取得（Docker環境で設定）
    voicevox_url = os.getenv('VOICEVOX_URL')
    if voicevox_url:
        return voicevox_url
    
    # Docker Compose環境では voicevox サービス名で接続
    try:
        # Docker環境かチェック
        if os.path.exists('/.dockerenv'):
            return "http://voicevox:50021"
    except:
        pass
    
    # WSL環境の場合はWindowsホストIPを取得
    try:
        with open('/etc/resolv.conf', 'r') as f:
            for line in f:
                if line.startswith('nameserver'):
                    host_ip = line.split()[1]
                    return f"http://{host_ip}:50021"
    except:
        pass
    
    # フォールバック: localhost
    return "http://localhost:50021"

def generate_voice_with_voicevox(text, speaker_id=10, output_path=None):
    """VOICEVOXを使用して音声を生成（雨晴はう: speaker_id=10）"""
    import requests
    import json
    import tempfile
    
    if output_path is None:
        output_path = tempfile.mktemp(suffix='.wav')
    
    # 環境に適したVOICEVOX URLを取得
    base_url = get_voicevox_url()
    
    try:
        # VOICEVOXエンジンの起動確認
        response = requests.get(f"{base_url}/speakers", timeout=5)
        if response.status_code != 200:
            raise Exception(f"VOICEVOXエンジンが起動していません (接続先: {base_url})")
        
        # 音響特徴量の生成
        query_response = requests.post(
            f"{base_url}/audio_query?text={text}&speaker={speaker_id}",
            timeout=10
        )
        query_response.raise_for_status()
        query_data = query_response.json()
        
        # 音声合成
        synthesis_response = requests.post(
            f"{base_url}/synthesis?speaker={speaker_id}",
            headers={"Content-Type": "application/json"},
            data=json.dumps(query_data),
            timeout=30
        )
        synthesis_response.raise_for_status()
        
        # 音声ファイルを保存
        with open(output_path, 'wb') as f:
            f.write(synthesis_response.content)
        
        return output_path
        
    except requests.exceptions.RequestException as e:
        raise Exception(f"VOICEVOXとの通信に失敗しました (接続先: {base_url}): {str(e)}")
    except Exception as e:
        raise Exception(f"音声生成に失敗しました: {str(e)}")

def add_multiple_voices_to_video(video_path, output_path, voices, original_volume=1.0):
    """動画に複数の音声を追加"""
    clip = VideoFileClip(video_path)
    
    # 元の音声を取得
    audio_clips = []
    if clip.audio is not None:
        original_audio = clip.audio
        if original_volume != 1.0:
            original_audio = original_audio.with_volume_scaled(original_volume)
        audio_clips.append(original_audio)
    
    # 各音声を生成して追加
    temp_voice_files = []
    try:
        for voice in voices:
            try:
                # 音声を生成
                voice_path = generate_voice_with_voicevox(voice['text'])
                temp_voice_files.append(voice_path)
                
                # 音声ファイルを読み込み
                voice_audio = AudioFileClip(voice_path)
                
                # 音量を調整
                if voice['volume'] != 1.0:
                    voice_audio = voice_audio.with_volume_scaled(voice['volume'])
                
                # 開始時間を設定（指定時間から再生開始）
                if voice['start_time'] > 0.0:
                    voice_audio = voice_audio.with_start(voice['start_time'])
                
                audio_clips.append(voice_audio)
                
            except Exception as e:
                st.warning(f"⚠️ 音声「{voice['text'][:20]}...」の生成をスキップしました: {str(e)}")
                continue
        
        # 全ての音声を合成
        if len(audio_clips) > 1:
            final_audio = CompositeAudioClip(audio_clips)
        elif len(audio_clips) == 1:
            final_audio = audio_clips[0]
        else:
            final_audio = None
        
        # 動画に音声を設定
        if final_audio is not None:
            final_clip = clip.with_audio(final_audio)
        else:
            final_clip = clip
        
        # 出力
        final_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        
        # リソースをクリーンアップ
        clip.close()
        final_clip.close()
        if final_audio is not None:
            final_audio.close()
        for audio in audio_clips:
            try:
                audio.close()
            except:
                pass
                
    finally:
        # 一時音声ファイルを削除
        for temp_file in temp_voice_files:
            try:
                os.unlink(temp_file)
            except:
                pass
    
    return output_path

def add_bgm_to_video(video_path, output_path, bgm_path=None, bgm_volume=0.5, original_volume=1.0, loop_bgm=True, bgm_start_time=0.0):
    """動画にBGMを追加"""
    clip = VideoFileClip(video_path)
    
    if bgm_path and os.path.exists(bgm_path):
        # BGMを読み込み
        bgm = AudioFileClip(bgm_path)
        
        # BGMの音量を調整
        bgm = bgm.with_fps(44100)
        if bgm_volume != 1.0:
            bgm = bgm.with_volume_scaled(bgm_volume)
        
        # BGMをループ再生するかどうか
        if loop_bgm and bgm.duration < clip.duration:
            # BGMをループして動画の長さに合わせる
            loops_needed = int(clip.duration / bgm.duration) + 1
            try:
                # ループ処理を行い、動画の長さに合わせてカット
                bgm = bgm.audio_loop(n=loops_needed).subclipped(0, clip.duration)
            except Exception as e:
                # ループ処理に失敗した場合は、単純に繰り返し再生
                bgm_list = []
                current_duration = 0
                while current_duration < clip.duration:
                    remaining_duration = clip.duration - current_duration
                    if remaining_duration >= bgm.duration:
                        bgm_list.append(bgm)
                        current_duration += bgm.duration
                    else:
                        bgm_list.append(bgm.subclipped(0, remaining_duration))
                        current_duration += remaining_duration
                
                if bgm_list:
                    bgm = concatenate_audioclips(bgm_list)
                else:
                    bgm = bgm.subclipped(0, min(bgm.duration, clip.duration))
        else:
            # BGMを動画の長さに合わせてカット
            bgm = bgm.subclipped(0, min(bgm.duration, clip.duration))
        
        # BGMの開始時間を適用
        if bgm_start_time > 0.0 and bgm_start_time < clip.duration:
            # 無音の音声を作成してBGMの前に追加
            silence_duration = min(bgm_start_time, clip.duration)
            remaining_duration = clip.duration - silence_duration
            
            if remaining_duration > 0:
                # BGMを残り時間に合わせてカット
                bgm = bgm.subclipped(0, min(bgm.duration, remaining_duration))
                # 無音とBGMを結合
                bgm = bgm.with_start(bgm_start_time)
            else:
                # 開始時間が動画の長さを超えている場合は無音
                bgm = None
        
        # 元の音声があるかチェック
        if clip.audio is not None:
            # 元の音声の音量を調整
            original_audio = clip.audio
            if original_volume != 1.0:
                original_audio = original_audio.with_volume_scaled(original_volume)
            # BGMと元の音声を合成
            if bgm is not None:
                final_audio = CompositeAudioClip([original_audio, bgm])
            else:
                final_audio = original_audio
        else:
            # 元の音声がない場合はBGMのみ
            if bgm is not None:
                final_audio = bgm
            else:
                final_audio = None
        
        # 動画に音声を設定
        if final_audio is not None:
            final_clip = clip.with_audio(final_audio)
        else:
            final_clip = clip
        
        # 出力
        final_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        
        # リソースをクリーンアップ
        clip.close()
        bgm.close()
        final_clip.close()
        if clip.audio is not None:
            original_audio.close()
        final_audio.close()
    else:
        # BGMがない場合は元の動画をそのままコピー
        clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        clip.close()
    
    return output_path

def add_text_to_video(video_path, output_path, telops, font_size=60):
    """動画に時間ベースのテキストオーバーレイを追加"""
    import cv2
    import numpy as np
    from PIL import Image, ImageDraw, ImageFont
    
    # OpenCVとPillowを使用してテキストを追加
    clip = VideoFileClip(video_path)
    
    def add_text_frame(get_frame, t):
        frame = get_frame(t)
        # numpy配列をPIL Imageに変換
        img = Image.fromarray(frame.astype('uint8'))
        draw = ImageDraw.Draw(img)
        
        # 日本語対応フォントを優先的に使用
        japanese_fonts = [
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "fonts/NotoSansJP-Regular.ttf",
            "NotoSansCJK-Regular.ttc"
        ]
        
        font = None
        for font_path in japanese_fonts:
            try:
                font = ImageFont.truetype(font_path, font_size)
                break
            except:
                continue
        
        if font is None:
            try:
                font = ImageFont.load_default(size=font_size)
            except:
                font = ImageFont.load_default()
        
        # 現在の時間に表示すべきテロップを描画
        for telop in telops:
            # 時間範囲内かチェック
            if telop['start_time'] <= t <= telop['end_time']:
                text = telop['text']
                position = telop['position']
                
                if not text:  # 空のテキストはスキップ
                    continue
                    
                # テキストサイズを取得
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
                
                # 位置を計算（中央寄りに調整）
                if position == "top":
                    x = (img.width - text_width) // 2
                    y = img.height // 4 - text_height // 2  # 中央寄りの上
                elif position == "bottom":
                    x = (img.width - text_width) // 2
                    y = img.height * 3 // 4 - text_height // 2  # 中央寄りの下
                else:  # center
                    x = (img.width - text_width) // 2
                    y = (img.height - text_height) // 2
                
                # テロップの色設定（デフォルトは白）
                text_color = telop.get('color', (255, 255, 255))
                
                # テキストを描画（影付き）
                draw.text((x+2, y+2), text, font=font, fill=(0, 0, 0))  # 影
                draw.text((x, y), text, font=font, fill=text_color)  # テキスト
        
        return np.array(img)
    
    # テキスト付きの動画を作成
    final_video = clip.transform(add_text_frame)
    
    # 出力（高画質設定）
    final_video.write_videofile(
        output_path, 
        codec='libx264', 
        audio_codec='aac',
        bitrate='8000k',  # 高ビットレート
        ffmpeg_params=['-crf', '18', '-preset', 'slow']  # 高画質・低圧縮
    )
    clip.close()
    final_video.close()
    
    return output_path

def combine_videos(video_paths, output_path):
    """複数の動画を結合する"""
    clips = []
    try:
        for video_path in video_paths:
            # ファイルの存在確認
            if not os.path.exists(video_path):
                raise FileNotFoundError(f"ファイルが見つかりません: {video_path}")
            
            clip = VideoFileClip(video_path)
            clips.append(clip)
        
        # 動画を結合
        final_clip = concatenate_videoclips(clips, method="compose")
        
        # 出力（高画質設定）
        final_clip.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            bitrate='8000k',
            ffmpeg_params=['-crf', '18', '-preset', 'slow']
        )
        
        return output_path
        
    finally:
        # リソースをクリーンアップ
        for clip in clips:
            try:
                clip.close()
            except:
                pass
        try:
            final_clip.close()
        except:
            pass

def extract_slides_from_pptx(pptx_file, output_dir):
    """PowerPointファイルをPNG画像に変換して抽出（LibreOffice使用）"""
    import subprocess
    import tempfile
    import os
    from PIL import Image
    
    try:
        print("DEBUG: Starting PowerPoint to PNG image conversion...")
        
        # 一時ファイルとディレクトリの作成
        temp_dir = tempfile.mkdtemp()
        temp_pptx_path = os.path.join(temp_dir, "presentation.pptx")
        
        # PowerPointファイルを一時保存
        if hasattr(pptx_file, 'getvalue'):
            pptx_data = pptx_file.getvalue()
        elif hasattr(pptx_file, 'read'):
            pptx_data = pptx_file.read()
        else:
            pptx_data = pptx_file
            
        with open(temp_pptx_path, 'wb') as f:
            f.write(pptx_data)
        
        print(f"DEBUG: Saved PowerPoint file to: {temp_pptx_path}")
        
        # LibreOfficeを使用してPowerPointを正確にPNG画像に変換
        libreoffice_cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'png',
            '--outdir', temp_dir,
            temp_pptx_path
        ]
        
        print(f"DEBUG: Running LibreOffice command: {' '.join(libreoffice_cmd)}")
        
        try:
            result = subprocess.run(libreoffice_cmd, capture_output=True, text=True, timeout=120)
            print(f"DEBUG: LibreOffice return code: {result.returncode}")
            print(f"DEBUG: LibreOffice stdout: {result.stdout}")
            print(f"DEBUG: LibreOffice stderr: {result.stderr}")
            
            # LibreOfficeは複数スライドを1つのPNGにしか変換しないため、常にフォールバック処理を使用
            print(f"DEBUG: LibreOffice only creates single PNG for multi-slide presentations, using enhanced method")
            raise Exception("LibreOffice single-PNG limitation, using enhanced slide creation")
                
        except subprocess.TimeoutExpired:
            print("DEBUG: LibreOffice conversion timed out, falling back to enhanced method")
            raise Exception("LibreOffice conversion timed out")
        except FileNotFoundError:
            print("DEBUG: LibreOffice not found, falling back to enhanced method")
            raise Exception("LibreOffice not found. Please install LibreOffice in the container.")
        
        # 変換された画像ファイルを探してリサイズ
        slide_images = []
        all_files = os.listdir(temp_dir)
        png_files = []
        
        # presentation-001.png, presentation-002.png のような形式を探す
        for filename in all_files:
            if filename.endswith('.png') and ('presentation' in filename or filename.startswith('Slide') or 'slide' in filename.lower()):
                png_files.append(filename)
        
        # ファイル名でソート
        png_files.sort()
        print(f"DEBUG: Directory contents: {all_files}")
        print(f"DEBUG: Found {len(png_files)} PNG slide files: {png_files}")
        
        if not png_files:
            # LibreOfficeが失敗した場合のフォールバック
            print("DEBUG: No PNG files found from LibreOffice, trying enhanced slide creation...")
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
            return create_enhanced_slides_with_png_save(pptx_file, output_dir)
        
        for i, png_file in enumerate(png_files):
            src_path = os.path.join(temp_dir, png_file)
            dst_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            
            # LibreOffice変換画像をそのまま1920x1080にリサイズ
            try:
                with Image.open(src_path) as img:
                    # 元の画像をそのまま保持し、適切にリサイズ
                    img = img.convert('RGB')  # RGBA→RGB変換（必要に応じて）
                    img_resized = img.resize((1920, 1080), Image.Resampling.LANCZOS)
                    img_resized.save(dst_path, 'PNG', quality=100)
                    slide_images.append(dst_path)
                    print(f"DEBUG: LibreOffice converted slide saved: {dst_path} (size: {img.size})")
            except Exception as img_error:
                print(f"DEBUG: Image processing failed for {png_file}: {img_error}")
                # エラー時はファイルをそのままコピー
                try:
                    import shutil
                    shutil.copy2(src_path, dst_path)
                    slide_images.append(dst_path)
                    print(f"DEBUG: Direct copy fallback: {dst_path}")
                except:
                    continue
        
        # 一時ファイルをクリーンアップ
        import shutil
        shutil.rmtree(temp_dir, ignore_errors=True)
        
        if slide_images:
            print(f"DEBUG: ✅ LibreOffice successfully converted {len(slide_images)} slides to PNG images")
        else:
            print(f"DEBUG: ❌ LibreOffice conversion produced no usable images")
        return slide_images, len(slide_images)
        
    except Exception as e:
        print(f"DEBUG: PowerPoint PNG conversion error: {e}")
        print(f"DEBUG: Falling back to image extraction method")
        # エラー時は画像抽出フォールバック手法を使用
        try:
            import shutil
            if 'temp_dir' in locals():
                shutil.rmtree(temp_dir, ignore_errors=True)
        except:
            pass
        return create_enhanced_slides_with_png_save(pptx_file, output_dir)

def create_image_based_slides(pptx_file, output_dir):
    """PowerPointから画像を直接抽出してスライド画像を作成"""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import io
    
    try:
        prs = Presentation(pptx_file)
        slide_images = []
        
        print(f"DEBUG: Processing {len(prs.slides)} slides with image extraction method")
        
        for i, slide in enumerate(prs.slides):
            # 1920x1080の背景画像を作成
            img = Image.new('RGB', (1920, 1080), 'white')
            
            # スライド内の画像を探して配置
            images_found = 0
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    try:
                        # PowerPointから画像データを抽出
                        image_part = shape.image.blob
                        ppt_image = Image.open(io.BytesIO(image_part))
                        
                        # 画像を1920x1080に収まるようにリサイズ
                        ppt_image.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                        
                        # 画像を中央に配置
                        x = (1920 - ppt_image.width) // 2
                        y = (1080 - ppt_image.height) // 2
                        img.paste(ppt_image, (x, y))
                        
                        images_found += 1
                        print(f"DEBUG: Slide {i+1} - Extracted image {images_found}")
                        break  # 最初の画像のみ使用
                        
                    except Exception as e:
                        print(f"DEBUG: Failed to extract image from slide {i+1}: {e}")
                        continue
            
            # 画像が見つからない場合は、テキストベースのフォールバック
            if images_found == 0:
                print(f"DEBUG: Slide {i+1} - No images found, using text fallback")
                img = create_text_slide(slide, i+1)
            
            # スライド画像を保存
            slide_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            img.save(slide_path, 'PNG', quality=95)
            slide_images.append(slide_path)
            print(f"DEBUG: Created image-based slide: {slide_path}")
        
        return slide_images, len(prs.slides)
        
    except Exception as e:
        print(f"DEBUG: Image extraction failed: {e}")
        # 最終フォールバック
        return create_enhanced_slides_with_png_save(pptx_file, output_dir)

def create_enhanced_slides_with_png_save(pptx_file, output_dir):
    """パワーポイントからスライドを作成してPNGで確実に保存"""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import io
    import os
    
    try:
        prs = Presentation(pptx_file)
        slide_images = []
        
        print(f"DEBUG: Processing {len(prs.slides)} slides with enhanced PNG creation method")
        print(f"DEBUG: PowerPoint file info - Total slides: {len(prs.slides)}")
        print(f"DEBUG: Slide dimensions: {prs.slide_width.emu if hasattr(prs.slide_width, 'emu') else 'N/A'} x {prs.slide_height.emu if hasattr(prs.slide_height, 'emu') else 'N/A'}")
        
        # 各スライドの詳細情報をログ出力（エラー回避版）
        for i, slide in enumerate(prs.slides):
            try:
                shape_count = len(slide.shapes) if hasattr(slide, 'shapes') else 0
                print(f"DEBUG: Slide {i+1} - Shapes: {shape_count}")
                if hasattr(slide, 'shapes') and shape_count > 0:
                    # python-pptxの既知のバグを回避するため、安全にshapeにアクセス
                    for j in range(min(3, shape_count)):
                        try:
                            shape = slide.shapes[j]
                            shape_type = type(shape).__name__
                            has_text = hasattr(shape, 'text') and bool(shape.text.strip())
                            print(f"DEBUG:   Shape {j+1}: {shape_type}, has_text: {has_text}")
                        except (AttributeError, IndexError) as shape_error:
                            print(f"DEBUG:   Shape {j+1}: Error accessing shape - {shape_error}")
            except Exception as slide_error:
                print(f"DEBUG: Slide {i+1} - Error analyzing slide: {slide_error}")
        
        
        # output_dirが存在しない場合は作成
        os.makedirs(output_dir, exist_ok=True)
        print(f"DEBUG: Output directory ensured: {output_dir}")
        
        for i, slide in enumerate(prs.slides):
            print(f"DEBUG: Creating slide {i+1}/{len(prs.slides)}...")
            
            # 最初にスライド内の画像を探して抽出を試みる
            slide_has_image = False
            img = Image.new('RGB', (1920, 1080), '#ffffff')
            
            # PowerPoint内の画像を探す（エラー回避版）
            try:
                if hasattr(slide, 'shapes'):
                    shape_count = len(slide.shapes)
                    for j in range(shape_count):
                        try:
                            shape = slide.shapes[j]
                            if hasattr(shape, 'image'):
                                # PowerPointから画像データを抽出
                                image_part = shape.image.blob
                                ppt_image = Image.open(io.BytesIO(image_part))
                                
                                # 画像を1920x1080にフィットさせる
                                ppt_image.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                                
                                # 中央に配置
                                x = (1920 - ppt_image.width) // 2
                                y = (1080 - ppt_image.height) // 2
                                img.paste(ppt_image, (x, y))
                                
                                slide_has_image = True
                                print(f"DEBUG: Slide {i+1} - Successfully extracted embedded image")
                                break  # 最初の画像のみ使用
                        except (AttributeError, IndexError) as shape_error:
                            print(f"DEBUG: Shape access error on slide {i+1}, shape {j}: {shape_error}")
                            continue
                        except Exception as e:
                            print(f"DEBUG: Failed to extract image from slide {i+1}, shape {j}: {e}")
                            continue
            except Exception as slide_error:
                print(f"DEBUG: Error processing shapes on slide {i+1}: {slide_error}")
            
            # 画像がない場合はテキストベースのスライドを作成
            if not slide_has_image:
                print(f"DEBUG: Slide {i+1} - No images found, creating text-based slide")
                img = create_enhanced_text_slide(slide, i+1, prs)
            
            # スライド画像を保存
            slide_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            
            try:
                img.save(slide_path, 'PNG', quality=95, optimize=True)
                
                # 保存が成功したか確認
                if os.path.exists(slide_path) and os.path.getsize(slide_path) > 0:
                    slide_images.append(slide_path)
                    print(f"DEBUG: ✅ Successfully saved slide: {slide_path} ({os.path.getsize(slide_path)} bytes)")
                else:
                    print(f"DEBUG: ❌ Failed to save slide: {slide_path}")
                    
            except Exception as save_error:
                print(f"DEBUG: Error saving slide {i+1}: {save_error}")
                continue
        
        print(f"DEBUG: ✅ Created {len(slide_images)} slide images in total")
        return slide_images, len(prs.slides)
        
    except Exception as e:
        print(f"DEBUG: Enhanced slide creation failed: {e}")
        import traceback
        print(f"DEBUG: Traceback: {traceback.format_exc()}")
        return [], 0

def create_enhanced_text_slide(slide, slide_num, presentation):
    """強化されたテキストベースのスライド画像を作成"""
    from PIL import Image, ImageDraw, ImageFont
    import textwrap
    
    # 高解像度で白い背景の画像を作成
    img = Image.new('RGB', (1920, 1080), '#ffffff')
    draw = ImageDraw.Draw(img)
    
    # 日本語フォントの読み込み
    try:
        font_title = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 64)
        font_content = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 36)
        font_small = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 28)
    except:
        # フォールバック
        font_title = ImageFont.load_default()
        font_content = ImageFont.load_default()
        font_small = ImageFont.load_default()
    
    # スライドからテキストを抽出（エラー回避版）
    texts_with_position = []
    try:
        if hasattr(slide, 'shapes'):
            shape_count = len(slide.shapes)
            for j in range(shape_count):
                try:
                    shape = slide.shapes[j]
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        texts_with_position.append({
                            'text': shape.text.strip(),
                            'top': getattr(shape, 'top', 0),
                            'left': getattr(shape, 'left', 0)
                        })
                except (AttributeError, IndexError) as shape_error:
                    print(f"DEBUG: Text extraction error for shape {j}: {shape_error}")
                    continue
                except Exception as e:
                    print(f"DEBUG: Unexpected error extracting text from shape {j}: {e}")
                    continue
    except Exception as slide_error:
        print(f"DEBUG: Error accessing slide shapes for text extraction: {slide_error}")
    
    # Y位置でソート
    texts_with_position.sort(key=lambda x: x['top'])
    
    # タイトルとコンテンツを分離
    if texts_with_position:
        title_text = texts_with_position[0]['text']
        content_texts = [item['text'] for item in texts_with_position[1:]]
    else:
        title_text = f"スライド {slide_num}"
        content_texts = []
    
    # 背景のグラデーション
    for y in range(80):
        color_val = 248 + int(y * 0.1)
        draw.rectangle([0, y, 1920, y+1], fill=(color_val, color_val, color_val+2))
    
    # タイトルを描画
    y_pos = 120
    wrapped_title = textwrap.fill(title_text, width=30)
    for line in wrapped_title.split('\n')[:2]:
        if line.strip():
            bbox = draw.textbbox((0, 0), line, font=font_title)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            x_pos = (1920 - text_width) // 2
            
            # タイトル背景
            padding = 20
            draw.rectangle([x_pos - padding, y_pos - 5, x_pos + text_width + padding, y_pos + text_height + 5], 
                          fill='#e8f4fd', outline='#2196f3', width=2)
            
            # タイトルテキスト（影付き）
            draw.text((x_pos + 1, y_pos + 1), line, fill='#333333', font=font_title)
            draw.text((x_pos, y_pos), line, fill='#1565c0', font=font_title)
            y_pos += text_height + 30
    
    # コンテンツを描画
    y_pos = max(y_pos + 50, 300)
    
    for content in content_texts[:5]:  # 最大5つのコンテンツ
        content_lines = content.replace('\n', '\n').split('\n')
        
        for line in content_lines[:10]:  # 行数制限
            if line.strip():
                wrapped_lines = textwrap.fill(line.strip(), width=50).split('\n')
                
                for wrapped_line in wrapped_lines:
                    if y_pos > 950:  # 画面下部に達した場合は停止
                        break
                        
                    # 箇条書きマーカー
                    marker_x = 150
                    draw.ellipse([marker_x - 8, y_pos + 12, marker_x, y_pos + 20], fill='#4caf50')
                    
                    # テキスト（影付き）
                    draw.text((marker_x + 21, y_pos + 1), wrapped_line, fill='#cccccc', font=font_content)
                    draw.text((marker_x + 20, y_pos), wrapped_line, fill='#333333', font=font_content)
                    
                    y_pos += 40
        
        y_pos += 20  # コンテンツ間の余白
    
    # フッター
    footer_text = f"スライド {slide_num} / {len(presentation.slides)}"
    footer_bbox = draw.textbbox((0, 0), footer_text, font=font_small)
    footer_width = footer_bbox[2] - footer_bbox[0]
    draw.rectangle([1920 - footer_width - 60, 1020, 1920 - 20, 1050], fill='#f5f5f5', outline='#ddd')
    draw.text((1920 - footer_width - 40, 1025), footer_text, fill='#666666', font=font_small)
    
    # 枠線
    draw.rectangle([10, 10, 1910, 1070], outline='#dddddd', width=2)
    
    return img

def create_text_slide(slide, slide_num):
    """単一スライドのテキストベース画像を作成"""
    from PIL import Image, ImageDraw, ImageFont
    import textwrap
    
    img = Image.new('RGB', (1920, 1080), '#f8f9fa')
    draw = ImageDraw.Draw(img)
    
    # フォント設定
    try:
        font_title = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 72)
        font_content = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 40)
    except:
        font_title = ImageFont.load_default()
        font_content = ImageFont.load_default()
    
    # タイトル部分
    title = f"スライド {slide_num}"
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            if len(shape.text.strip()) < 50:
                title = shape.text.strip()
                break
    
    # タイトル描画
    title_bbox = draw.textbbox((0, 0), title, font=font_title)
    title_x = (1920 - (title_bbox[2] - title_bbox[0])) // 2
    draw.text((title_x, 200), title, fill='#1565c0', font=font_title)
    
    # コンテンツ描画
    y_pos = 400
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if text != title:  # タイトル以外のテキスト
                wrapped_lines = textwrap.fill(text, width=45).split('\n')
                for line in wrapped_lines[:8]:  # 最大8行
                    if line.strip():
                        line_bbox = draw.textbbox((0, 0), line, font=font_content)
                        line_x = (1920 - (line_bbox[2] - line_bbox[0])) // 2
                        draw.text((line_x, y_pos), line, fill='#424242', font=font_content)
                        y_pos += 60
    
    return img

def create_fallback_slides(pptx_file, output_dir):
    """改良されたスライド画像作成（日本語対応・レイアウト改善）"""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import textwrap
    
    try:
        prs = Presentation(pptx_file)
        slide_images = []
        
        print(f"DEBUG: Processing {len(prs.slides)} slides with enhanced PNG fallback method")
        
        for i, slide in enumerate(prs.slides):
            # 高品質なスライド画像を作成
            img = Image.new('RGB', (1920, 1080), '#f8f9fa')  # 少し灰色がかった白
            draw = ImageDraw.Draw(img)
            
            # 日本語対応フォントを読み込み
            try:
                font_title = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 72)
                font_subtitle = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 48)
                font_content = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 40)
                print(f"DEBUG: Successfully loaded Japanese fonts for slide {i+1}")
            except Exception as e:
                print(f"DEBUG: Font loading failed for slide {i+1}: {e}")
                font_title = ImageFont.load_default()
                font_subtitle = ImageFont.load_default()
                font_content = ImageFont.load_default()
            
            # スライドコンテンツを分析・抽出
            title_text = ""
            content_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # タイトル判定（最初の図形または短いテキスト）
                    if not title_text and (len(text) < 50 or shape.top < prs.slide_height * 0.3):
                        title_text = text
                    else:
                        content_texts.append(text)
            
            # デフォルトタイトル設定
            if not title_text:
                title_text = f"スライド {i+1}"
            
            print(f"DEBUG: Slide {i+1} - Title: '{title_text[:30]}...' Content items: {len(content_texts)}")
            
            # タイトル描画（上部中央）
            title_lines = title_text.split('\\n') if '\\n' in title_text else [title_text]
            y_pos = 120
            
            for title_line in title_lines[:2]:  # 最大2行
                if title_line.strip():
                    # 長いタイトルを折り返し
                    wrapped_title = textwrap.fill(title_line, width=25)
                    for wrapped_line in wrapped_title.split('\\n'):
                        title_bbox = draw.textbbox((0, 0), wrapped_line, font=font_title)
                        title_x = (1920 - (title_bbox[2] - title_bbox[0])) // 2
                        
                        # タイトル背景
                        padding = 20
                        draw.rectangle([
                            title_x - padding, y_pos - padding,
                            title_x + (title_bbox[2] - title_bbox[0]) + padding, 
                            y_pos + (title_bbox[3] - title_bbox[1]) + padding
                        ], fill='#e3f2fd', outline='#1976d2', width=2)
                        
                        # タイトルテキスト
                        draw.text((title_x, y_pos), wrapped_line, fill='#1565c0', font=font_title)
                        y_pos += 90
            
            # コンテンツ描画（中央部分）
            y_pos = 350
            max_content_lines = 12
            line_count = 0
            
            for content in content_texts:
                if line_count >= max_content_lines:
                    break
                    
                content_lines = content.split('\\n')
                for line in content_lines:
                    if line_count >= max_content_lines:
                        break
                        
                    if line.strip():
                        # 長い行を折り返し
                        wrapped_lines = textwrap.fill(line.strip(), width=45).split('\\n')
                        
                        for wrapped_line in wrapped_lines:
                            if line_count >= max_content_lines:
                                break
                                
                            # 内容を左揃えで配置
                            margin_left = 160
                            
                            # 箇条書きマーカー
                            draw.ellipse([margin_left - 20, y_pos + 15, margin_left - 10, y_pos + 25], fill='#4caf50')
                            
                            # コンテンツテキスト
                            draw.text((margin_left, y_pos), wrapped_line, fill='#424242', font=font_content)
                            
                            y_pos += 55
                            line_count += 1
                
                # セクション間のスペース
                if content_texts.index(content) < len(content_texts) - 1:
                    y_pos += 20
            
            # フッター（スライド番号）
            footer_text = f"{i+1} / {len(prs.slides)}"
            footer_bbox = draw.textbbox((0, 0), footer_text, font=font_content)
            draw.text((1920 - (footer_bbox[2] - footer_bbox[0]) - 50, 1020), 
                     footer_text, fill='#757575', font=font_content)
            
            # 装飾的な枠線
            draw.rectangle([30, 30, 1890, 1050], outline='#e0e0e0', width=3)
            
            # 画像を保存
            slide_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            img.save(slide_path, 'PNG', quality=95)
            slide_images.append(slide_path)
            print(f"DEBUG: Created enhanced slide image: {slide_path}")
        
        return slide_images, len(prs.slides)
        
    except Exception as e:
        print(f"DEBUG: Enhanced slide creation failed: {e}")
        raise Exception(f"Enhanced slide creation failed: {str(e)}")

def create_slide_video(slide_images, durations, bgm_path, lyrics, output_path):
    """スライド画像から動画を作成（FFmpeg直接使用）"""
    import subprocess
    import tempfile
    
    print(f"DEBUG: create_slide_video called with {len(slide_images)} slides")
    print(f"DEBUG: Using FFmpeg-based approach, no MoviePy")
    
    try:
        # 一時ファイルのリスト
        temp_files = []
        
        # 1. 各スライドから個別の動画を作成
        slide_videos = []
        for i, (slide_path, duration) in enumerate(zip(slide_images, durations)):
            temp_video = tempfile.mktemp(suffix=f'_slide_{i}.mp4')
            temp_files.append(temp_video)
            
            # FFmpegでスライド画像から動画を作成
            ffmpeg_cmd = [
                'ffmpeg', '-y',  # -y は上書き確認
                '-loop', '1',  # 画像をループ
                '-i', slide_path,  # 入力画像
                '-t', str(duration),  # 持続時間
                '-r', '24',  # FPSを明示的に設定
                '-c:v', 'libx264',
                '-pix_fmt', 'yuv420p',
                '-crf', '18',
                temp_video
            ]
            
            result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
            if result.returncode != 0:
                raise Exception(f"FFmpegエラー (slide {i}): {result.stderr}")
            
            slide_videos.append(temp_video)
        
        # 2. スライド動画を結合
        concat_file = tempfile.mktemp(suffix='_concat.txt')
        temp_files.append(concat_file)
        
        with open(concat_file, 'w') as f:
            for video in slide_videos:
                f.write(f"file '{video}'\n")
        
        temp_combined = tempfile.mktemp(suffix='_combined.mp4')
        temp_files.append(temp_combined)
        
        # FFmpegで動画結合
        concat_cmd = [
            'ffmpeg', '-y',
            '-f', 'concat',
            '-safe', '0',
            '-i', concat_file,
            '-c', 'copy',
            temp_combined
        ]
        
        result = subprocess.run(concat_cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise Exception(f"FFmpeg結合エラー: {result.stderr}")
        
        # 3. BGMを追加（オプション）
        current_video = temp_combined
        if bgm_path and os.path.exists(bgm_path):
            temp_with_bgm = tempfile.mktemp(suffix='_with_bgm.mp4')
            temp_files.append(temp_with_bgm)
            
            # BGMを追加（ループして動画の長さに合わせる）
            bgm_cmd = [
                'ffmpeg', '-y',
                '-i', current_video,  # 動画入力
                '-stream_loop', '-1',  # BGMをループ
                '-i', bgm_path,  # BGM入力
                '-shortest',  # 短い方に合わせる
                '-c:v', 'copy',  # 動画はコピー
                '-c:a', 'aac',  # 音声はAAC
                '-filter:a', 'volume=0.3',  # BGM音量を下げる
                temp_with_bgm
            ]
            
            result = subprocess.run(bgm_cmd, capture_output=True, text=True)
            if result.returncode == 0:
                current_video = temp_with_bgm
            else:
                print(f"BGM追加でエラー: {result.stderr}")
                # BGMエラーは無視して続行
        
        # 4. 歌詞を追加（オプション）
        if lyrics:
            temp_with_lyrics = tempfile.mktemp(suffix='_with_lyrics.mp4')
            temp_files.append(temp_with_lyrics)
            
            # 歌詞フィルターを作成
            drawtext_filters = []
            for lyric in lyrics:
                # 特殊文字をエスケープ
                text = lyric['text'].replace(":", "\\:").replace("'", "\\'")
                filter_str = f"drawtext=text='{text}':fontsize=40:fontcolor=white:x=(w-text_w)/2:y=h-100:enable='between(t,{lyric['start_time']},{lyric['end_time']})'"
                drawtext_filters.append(filter_str)
            
            if drawtext_filters:
                filter_complex = ','.join(drawtext_filters)
                
                lyrics_cmd = [
                    'ffmpeg', '-y',
                    '-i', current_video,
                    '-vf', filter_complex,
                    '-c:a', 'copy',  # 音声はコピー
                    temp_with_lyrics
                ]
                
                result = subprocess.run(lyrics_cmd, capture_output=True, text=True)
                if result.returncode == 0:
                    current_video = temp_with_lyrics
                else:
                    print(f"歌詞追加でエラー: {result.stderr}")
                    # 歌詞エラーは無視して続行
        
        # 5. 最終出力
        final_cmd = [
            'ffmpeg', '-y',
            '-i', current_video,
            '-c:v', 'libx264',
            '-c:a', 'aac',
            '-crf', '18',
            '-preset', 'slow',
            output_path
        ]
        
        result = subprocess.run(final_cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise Exception(f"FFmpeg最終出力エラー: {result.stderr}")
        
        # 一時ファイルをクリーンアップ
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except:
                pass
        
        return output_path
        
    except Exception as e:
        # エラー時も一時ファイルをクリーンアップ
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except:
                pass
        raise Exception(f"動画作成に失敗しました: {str(e)}")

def add_lyrics_to_video(video_clip, lyrics):
    """動画に歌詞を追加（非使用・代替手法あり）"""
    # この関数は現在使用されていません
    # create_slide_video内でTextClipを使用して歌詞を追加しています
    return video_clip

# メインインターface
if tool == "ショート動画変換":
    uploaded_file = st.file_uploader(
        "動画ファイルをアップロードしてください",
        type=['mp4', 'avi', 'mov', 'mkv'],
        help="MP4, AVI, MOV, MKV形式の動画ファイルをサポートしています"
    )
elif tool == "動画結合":
    uploaded_files = st.file_uploader(
        "結合する動画ファイルを複数選択してください",
        type=['mp4', 'avi', 'mov', 'mkv'],
        accept_multiple_files=True,
        help="MP4, AVI, MOV, MKV形式の動画ファイルをサポートしています"
    )
# スライド動画作成ツールには動画アップローダーは不要（PowerPointファイル用の専用UI）

if tool == "ショート動画変換" and uploaded_file is not None:
    # 一時ファイルに保存
    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp_file:
        tmp_file.write(uploaded_file.read())
        input_video_path = tmp_file.name
    
    # 動画情報を表示
    try:
        clip = VideoFileClip(input_video_path)
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("解像度", f"{clip.w}x{clip.h}")
        with col2:
            st.metric("時間", f"{clip.duration:.1f}秒")
        with col3:
            st.metric("FPS", f"{clip.fps:.1f}")
        
        clip.close()
        
        # オプション設定
        st.subheader("設定オプション")
        
        # 動画トリミング設定
        st.subheader("✂️ 動画トリミング設定")
        trim_video = st.checkbox("動画の一部分だけを抜き出す", help="元動画の指定した時間範囲のみを使用します")
        
        if trim_video:
            st.markdown("**📌 抜き出したい時間範囲を指定してください**")
            
            # セッション状態でトリミング値を管理
            if 'trim_start_time' not in st.session_state:
                st.session_state.trim_start_time = 0.0
            if 'trim_end_time' not in st.session_state:
                st.session_state.trim_end_time = min(30.0, float(clip.duration))
            
            # クイック選択プリセット
            st.markdown("**🚀 クイック選択**")
            preset_col1, preset_col2, preset_col3, preset_col4 = st.columns(4)
            
            with preset_col1:
                if st.button("📺 最初の30秒", help="動画の最初から30秒を抜き出し"):
                    st.session_state.trim_start_time = 0.0
                    st.session_state.trim_end_time = min(30.0, float(clip.duration))
            with preset_col2:
                if st.button("🎯 中央30秒", help="動画の中央部分30秒を抜き出し"):
                    mid_point = clip.duration / 2
                    st.session_state.trim_start_time = max(0.0, mid_point - 15.0)
                    st.session_state.trim_end_time = min(float(clip.duration), mid_point + 15.0)
            with preset_col3:
                if st.button("🏁 最後の30秒", help="動画の最後から30秒を抜き出し"):
                    st.session_state.trim_start_time = max(0.0, float(clip.duration) - 30.0)
                    st.session_state.trim_end_time = float(clip.duration)
            with preset_col4:
                if st.button("🎬 全体使用", help="動画全体を使用"):
                    st.session_state.trim_start_time = 0.0
                    st.session_state.trim_end_time = float(clip.duration)
            
            st.markdown("**⚙️ 詳細設定**")
            
            # スライダー形式での時間範囲指定
            col1, col2, col3 = st.columns([1, 2, 1])
            
            with col1:
                start_time = st.number_input(
                    "開始時間（秒）", 
                    min_value=0.0, 
                    max_value=float(clip.duration - 0.1), 
                    value=st.session_state.trim_start_time, 
                    step=0.1,
                    format="%.1f",
                    help="この時間から動画を開始します",
                    key="start_time_input"
                )
                st.session_state.trim_start_time = start_time
            
            with col2:
                # 時間範囲スライダー
                time_range = st.slider(
                    "時間範囲を視覚的に選択",
                    min_value=0.0,
                    max_value=float(clip.duration),
                    value=(st.session_state.trim_start_time, st.session_state.trim_end_time),
                    step=0.1,
                    format="%.1fs",
                    key="time_range_slider"
                )
                st.session_state.trim_start_time = time_range[0]
                st.session_state.trim_end_time = time_range[1]
                start_time = time_range[0]
                end_time = time_range[1]
            
            with col3:
                end_time = st.number_input(
                    "終了時間（秒）", 
                    min_value=st.session_state.trim_start_time + 0.1, 
                    max_value=float(clip.duration), 
                    value=st.session_state.trim_end_time, 
                    step=0.1,
                    format="%.1f",
                    help="この時間で動画を終了します",
                    key="end_time_input"
                )
                st.session_state.trim_end_time = end_time
            
            # トリミング情報の表示
            trimmed_duration = end_time - start_time
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.metric("🎬 元動画の長さ", f"{clip.duration:.1f}秒")
            with col2:
                st.metric("✂️ 抜き出す範囲", f"{start_time:.1f}s - {end_time:.1f}s")
            with col3:
                st.metric("⏱️ 抜き出し後の長さ", f"{trimmed_duration:.1f}秒")
            
            # プレビュー情報
            if trimmed_duration > 0:
                percentage = (trimmed_duration / clip.duration) * 100
                st.success(f"✅ 元動画の **{percentage:.1f}%** を抜き出します（{start_time:.1f}秒 ～ {end_time:.1f}秒）")
                
                # YouTubeショートの推奨時間チェック
                if trimmed_duration > 60:
                    st.warning("⚠️ YouTubeショートは60秒以内が推奨です")
                elif trimmed_duration < 5:
                    st.warning("⚠️ 動画が短すぎる可能性があります（5秒未満）")
            else:
                st.error("❌ 終了時間は開始時間より後に設定してください")
        
        # スケール倍率設定
        st.subheader("📏 拡大倍率設定")
        scale_factor = st.slider(
            "動画の拡大倍率",
            min_value=0.5,
            max_value=5.0,
            value=1.0,
            step=0.1,
            help="1.0 = 原寸大、3.0 = 300%拡大。大きくするほどズームインされます。"
        )
        
        # 倍率の説明表示
        if scale_factor < 1.0:
            st.info(f"📉 縮小表示: {scale_factor*100:.0f}% (より多くの映像が見えます)")
        elif scale_factor == 1.0:
            st.info("📊 原寸大表示: 100% (元の動画のまま)")
        else:
            st.info(f"📈 拡大表示: {scale_factor*100:.0f}% (ズームイン効果)")
        
        # テキストオーバーレイ設定
        add_text = st.checkbox("テキストを追加する")
        
        if add_text:
            st.subheader("テロップ設定")
            
            # セッション状態でテロップリストを管理
            if 'telops' not in st.session_state:
                st.session_state.telops = []
            
            font_size = st.slider("フォントサイズ", 30, 120, 60)
            
            # フォント色の選択
            st.subheader("カラー設定")
            color_options = {
                "白": (255, 255, 255),
                "黒": (0, 0, 0),
                "赤": (255, 0, 0),
                "青": (0, 0, 255),
                "緑": (0, 255, 0),
                "黄": (255, 255, 0),
                "マゼンタ": (255, 0, 255),
                "シアン": (0, 255, 255),
                "オレンジ": (255, 165, 0)
            }
            selected_color_name = st.selectbox("フォント色", list(color_options.keys()), index=0)
            selected_color = color_options[selected_color_name]
            
            # 新しいテロップを追加するフォーム
            with st.expander("新しいテロップを追加", expanded=True):
                col1, col2, col3, col4, col5, col6 = st.columns([2, 1, 1, 1, 1, 1])
                
                with col1:
                    new_text = st.text_input("テロップテキスト", key="new_telop_text")
                with col2:
                    new_position = st.selectbox("位置", ["top", "center", "bottom"], key="new_telop_position")
                with col3:
                    new_start = st.number_input("開始(秒)", min_value=0, value=0, step=1, key="new_telop_start")
                with col4:
                    new_end = st.number_input("終了(秒)", min_value=1, value=5, step=1, key="new_telop_end")
                with col5:
                    new_color_name = st.selectbox("色", list(color_options.keys()), index=0, key="new_telop_color")
                with col6:
                    if st.button("追加", type="primary"):
                        if new_text.strip():
                            new_telop = {
                                "text": new_text,
                                "position": new_position,
                                "start_time": new_start,
                                "end_time": new_end,
                                "color": color_options[new_color_name]
                            }
                            st.session_state.telops.append(new_telop)
                            st.rerun()
            
            # 既存のテロップを表示・編集
            if st.session_state.telops:
                st.subheader("登録済みテロップ")
                
                for i, telop in enumerate(st.session_state.telops):
                    with st.container():
                        col1, col2, col3, col4, col5, col6, col7 = st.columns([2, 1, 1, 1, 1, 1, 0.5])
                        
                        with col1:
                            st.text_input(f"テキスト {i+1}", value=telop["text"], key=f"telop_text_{i}", disabled=True)
                        with col2:
                            st.text(telop["position"])
                        with col3:
                            st.text(f"{telop['start_time']}秒")
                        with col4:
                            st.text(f"{telop['end_time']}秒")
                        with col5:
                            duration = telop['end_time'] - telop['start_time']
                            st.text(f"{duration:.1f}秒間")
                        with col6:
                            # 色を表示
                            color = telop.get('color', (255, 255, 255))
                            color_name = next((name for name, rgb in color_options.items() if rgb == color), "白")
                            st.text(color_name)
                        with col7:
                            if st.button("🗑️", key=f"delete_telop_{i}", help="削除"):
                                st.session_state.telops.pop(i)
                                st.rerun()
                
                if st.button("全てのテロップをクリア", type="secondary"):
                    st.session_state.telops = []
                    st.rerun()
        
        # 音声合成設定（雨晴はう）
        st.subheader("🎤 音声合成設定（雨晴はう）")
        
        # VOICEVOXの設定情報を表示
        st.info("""
        💡 **VOICEVOX設定について (WSL環境)**:
        1. Windows側でVOICEVOXを起動してください
        2. VOICEVOXの設定で「外部からのアクセスを許可」を有効にしてください
        3. Windowsファイアウォールで port 50021 を開放してください
        """)
        
        add_voice = st.checkbox("雨晴はうの音声を追加する")
        
        if add_voice:
            # セッション状態で音声リストを管理
            if 'voices' not in st.session_state:
                st.session_state.voices = []
            
            # 新しい音声を追加するフォーム
            with st.expander("新しい音声を追加", expanded=True):
                col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                
                with col1:
                    new_voice_text = st.text_area(
                        "読み上げテキスト",
                        placeholder="雨晴はうに読み上げてもらいたいテキストを入力してください",
                        key="new_voice_text",
                        height=100
                    )
                with col2:
                    new_voice_start = st.number_input("開始時間(秒)", min_value=0, value=0, step=1, key="new_voice_start")
                with col3:
                    new_voice_volume = st.slider("音量", 0.0, 1.0, 0.8, 0.1, key="new_voice_volume")
                with col4:
                    st.write("") # スペース調整
                    if st.button("🔊 プレビュー", key="preview_voice"):
                        if new_voice_text.strip():
                            try:
                                with st.spinner("音声を生成中..."):
                                    voice_path = generate_voice_with_voicevox(new_voice_text)
                                    st.audio(voice_path)
                                    os.unlink(voice_path)
                                    st.success("✅ 音声生成成功！")
                            except Exception as e:
                                st.error(f"❌ 音声生成エラー: {str(e)}")
                    
                    if st.button("追加", type="primary", key="add_voice"):
                        if new_voice_text.strip():
                            new_voice = {
                                "text": new_voice_text,
                                "start_time": new_voice_start,
                                "volume": new_voice_volume
                            }
                            st.session_state.voices.append(new_voice)
                            st.rerun()
            
            # 既存の音声を表示・編集
            if st.session_state.voices:
                st.subheader("登録済み音声")
                
                for i, voice in enumerate(st.session_state.voices):
                    with st.container():
                        col1, col2, col3, col4 = st.columns([3, 1, 1, 0.5])
                        
                        with col1:
                            # テキストを1行で表示（長い場合は省略）
                            display_text = voice["text"][:50] + "..." if len(voice["text"]) > 50 else voice["text"]
                            st.text_input(f"音声 {i+1}", value=display_text, key=f"voice_text_{i}", disabled=True)
                        with col2:
                            st.text(f"{voice['start_time']}秒から")
                        with col3:
                            st.text(f"音量: {voice['volume']}")
                        with col4:
                            if st.button("🗑️", key=f"delete_voice_{i}", help="削除"):
                                st.session_state.voices.pop(i)
                                st.rerun()
                
                if st.button("全ての音声をクリア", type="secondary"):
                    st.session_state.voices = []
                    st.rerun()
        
        # BGM設定
        st.subheader("🎵 BGM設定")
        add_bgm = st.checkbox("BGMを追加する")
        
        if add_bgm:
            bgm_file = st.file_uploader(
                "BGMファイルをアップロード",
                type=['mp3', 'wav', 'aac', 'm4a', 'ogg'],
                help="MP3, WAV, AAC, M4A, OGG形式の音声ファイルをサポートしています"
            )
            
            col1, col2 = st.columns(2)
            with col1:
                bgm_volume = st.slider("BGM音量", 0.0, 1.0, 0.3, 0.1)
            with col2:
                original_volume = st.slider("元音声音量", 0.0, 1.0, 0.7, 0.1)
            
            loop_bgm = st.checkbox("BGMをループ再生", value=True)
            
            if bgm_file:
                st.audio(bgm_file)
        
        # 変換ボタン
        if st.button("ショート動画に変換", type="primary"):
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                # BGMファイルを一時保存
                bgm_path = None
                if add_bgm and bgm_file is not None:
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as tmp_bgm:
                        tmp_bgm.write(bgm_file.read())
                        bgm_path = tmp_bgm.name
                
                # Step 1: 動画をショート形式にリサイズ
                status_text.text("動画をリサイズ中...")
                progress_bar.progress(20)
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='_resized.mp4') as tmp_resized:
                    resized_video_path = tmp_resized.name
                
                # トリミングパラメータを設定
                trim_start = start_time if trim_video else None
                trim_end = end_time if trim_video else None
                
                resize_video_to_shorts(input_video_path, resized_video_path, scale_factor, trim_start, trim_end)
                progress_bar.progress(40)
                
                # Step 2: テキストを追加（オプション）
                current_video_path = resized_video_path
                if add_text and st.session_state.telops:
                    status_text.text("テキストを追加中...")
                    progress_bar.progress(60)
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix='_with_text.mp4') as tmp_text:
                        text_video_path = tmp_text.name
                    
                    add_text_to_video(current_video_path, text_video_path, st.session_state.telops, font_size)
                    os.unlink(current_video_path)
                    current_video_path = text_video_path
                
                # Step 3: 音声合成を追加（オプション）
                if add_voice and st.session_state.voices:
                    status_text.text("雨晴はうの音声を生成・追加中...")
                    progress_bar.progress(60)
                    
                    try:
                        # 複数音声を動画に追加
                        with tempfile.NamedTemporaryFile(delete=False, suffix='_with_voices.mp4') as tmp_voice:
                            voice_video_path = tmp_voice.name
                        
                        add_multiple_voices_to_video(
                            current_video_path,
                            voice_video_path,
                            st.session_state.voices,
                            1.0  # 元音声音量
                        )
                        
                        os.unlink(current_video_path)
                        current_video_path = voice_video_path
                        
                    except Exception as e:
                        st.warning(f"⚠️ 音声合成をスキップしました: {str(e)}")
                
                # Step 4: BGMを追加（オプション）
                if add_bgm and bgm_path:
                    status_text.text("BGMを追加中...")
                    progress_bar.progress(80)
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix='_final.mp4') as tmp_final:
                        final_video_path = tmp_final.name
                    
                    add_bgm_to_video(
                        current_video_path, 
                        final_video_path, 
                        bgm_path, 
                        bgm_volume, 
                        original_volume, 
                        loop_bgm
                    )
                    os.unlink(current_video_path)
                    os.unlink(bgm_path)
                else:
                    final_video_path = current_video_path
                
                progress_bar.progress(100)
                status_text.text("変換完了！")
                
                # プレビュー表示
                st.subheader("📹 プレビュー")
                st.video(final_video_path, width=400)
                
                # ダウンロードボタンを表示
                with open(final_video_path, 'rb') as file:
                    st.download_button(
                        label="📱 ショート動画をダウンロード",
                        data=file.read(),
                        file_name=f"shorts_{uploaded_file.name}",
                        mime="video/mp4"
                    )
                
                # 一時ファイルをクリーンアップ
                os.unlink(input_video_path)
                os.unlink(final_video_path)
                
                st.success("✅ 変換が完了しました！ダウンロードボタンをクリックして保存してください。")
                
            except Exception as e:
                st.error(f"❌ エラーが発生しました: {str(e)}")
                # エラー時のクリーンアップ
                try:
                    os.unlink(input_video_path)
                except:
                    pass
    
    except Exception as e:
        st.error(f"❌ 動画ファイルの読み込みに失敗しました: {str(e)}")
        os.unlink(input_video_path)

elif tool == "動画結合" and uploaded_files:
    if len(uploaded_files) < 2:
        st.warning("⚠️ 2つ以上の動画ファイルを選択してください。")
    else:
        st.success(f"✅ {len(uploaded_files)}個の動画ファイルが選択されました。")
        
        # 動画情報を表示
        st.subheader("📹 選択された動画")
        total_duration = 0
        
        for i, file in enumerate(uploaded_files):
            # ファイルを一時保存
            file.seek(0)  # ファイルポインタを先頭に戻す
            temp_path = tempfile.mktemp(suffix=f'_preview_{i}.mp4')
            
            try:
                with open(temp_path, 'wb') as tmp_file:
                    file_bytes = file.getvalue()
                    tmp_file.write(file_bytes)
                    tmp_file.flush()
                    os.fsync(tmp_file.fileno())
                
                # ファイルサイズを確認
                if os.path.getsize(temp_path) == 0:
                    st.error(f"❌ {file.name} のサイズが0です")
                    continue
                
            except Exception as e:
                st.error(f"❌ {file.name}の保存に失敗しました: {str(e)}")
                continue
            
            try:
                clip = VideoFileClip(temp_path)
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    st.text(f"{i+1}. {file.name}")
                with col2:
                    st.text(f"{clip.w}x{clip.h}")
                with col3:
                    st.text(f"{clip.duration:.1f}秒")
                with col4:
                    st.text(f"{clip.fps:.1f} FPS")
                
                total_duration += clip.duration
                clip.close()
                os.unlink(temp_path)
            except Exception as e:
                st.error(f"❌ {file.name}の読み込みに失敗しました: {str(e)}")
                try:
                    os.unlink(temp_path)
                except:
                    pass
        
        st.info(f"📊 結合後の総時間: {total_duration:.1f}秒")
        
        # 結合ボタン
        if st.button("動画を結合", type="primary"):
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                status_text.text("動画を結合中...")
                progress_bar.progress(20)
                
                # 一時ファイルに保存
                temp_paths = []
                for i, file in enumerate(uploaded_files):
                    file.seek(0)  # ファイルポインタを先頭に戻す
                    
                    # 一意のファイル名を生成
                    temp_path = tempfile.mktemp(suffix=f'_video_{i}.mp4')
                    
                    try:
                        with open(temp_path, 'wb') as tmp_file:
                            file_bytes = file.getvalue()
                            tmp_file.write(file_bytes)
                            tmp_file.flush()
                            os.fsync(tmp_file.fileno())
                        
                        # ファイルサイズを確認
                        if os.path.getsize(temp_path) == 0:
                            raise ValueError(f"ファイル {file.name} のサイズが0です")
                        
                        temp_paths.append(temp_path)
                        
                    except Exception as e:
                        st.error(f"❌ {file.name}の保存に失敗しました: {str(e)}")
                        # 失敗したファイルがあればクリーンアップして終了
                        for path in temp_paths:
                            try:
                                os.unlink(path)
                            except:
                                pass
                        raise
                
                progress_bar.progress(50)
                
                # 動画を結合
                with tempfile.NamedTemporaryFile(delete=False, suffix='_combined.mp4') as tmp_output:
                    output_path = tmp_output.name
                
                combine_videos(temp_paths, output_path)
                progress_bar.progress(100)
                status_text.text("結合完了！")
                
                # プレビュー表示
                st.subheader("📹 結合された動画")
                st.video(output_path, width=400)
                
                # ダウンロードボタンを表示
                with open(output_path, 'rb') as file:
                    st.download_button(
                        label="📱 結合動画をダウンロード",
                        data=file.read(),
                        file_name=f"combined_{len(uploaded_files)}_videos.mp4",
                        mime="video/mp4"
                    )
                
                # 一時ファイルをクリーンアップ
                for temp_path in temp_paths:
                    os.unlink(temp_path)
                os.unlink(output_path)
                
                st.success("✅ 結合が完了しました！ダウンロードボタンをクリックして保存してください。")
                
            except Exception as e:
                st.error(f"❌ エラーが発生しました: {str(e)}")
                # エラー時のクリーンアップ
                for temp_path in temp_paths:
                    try:
                        os.unlink(temp_path)
                    except:
                        pass

elif tool == "スライド動画作成":
    # PowerPointファイルのアップロード
    st.subheader("📂 ファイルアップロード")
    col1, col2 = st.columns(2)
    
    with col1:
        pptx_file = st.file_uploader(
            "PowerPointファイル",
            type=['pptx', 'ppt'],
            help="PowerPoint形式のプレゼンテーションファイルをアップロードしてください"
        )
    
    with col2:
        bgm_file = st.file_uploader(
            "BGMファイル",
            type=['mp3', 'wav', 'aac', 'm4a', 'ogg'],
            help="背景音楽ファイルをアップロードしてください"
        )
    
    if pptx_file is not None:
        # PowerPointファイルを処理
        try:
            # 一時ディレクトリを作成
            with tempfile.TemporaryDirectory() as temp_dir:
                # スライド画像を抽出 (PowerPoint → PNG画像変換)
                st.info("🔄 PowerPointをPNG画像に変換中...")
                with st.spinner("変換処理中..."):
                    slide_images, slide_count = extract_slides_from_pptx(pptx_file, temp_dir)
                
                # デバッグ情報をStreamlitに表示
                st.write(f"**デバッグ情報:**")
                st.write(f"- PowerPointスライド数: {slide_count}")
                st.write(f"- 生成されたPNG画像数: {len(slide_images)}")
                st.write(f"- PNG画像パス: {slide_images[:3]}{'...' if len(slide_images) > 3 else ''}")
                
                # ログにも出力
                print(f"DEBUG UI: PowerPoint slides: {slide_count}, Generated PNGs: {len(slide_images)}")
                for i, img_path in enumerate(slide_images[:5]):
                    print(f"DEBUG UI: Slide {i+1}: {img_path} (exists: {os.path.exists(img_path)})")
                
                if slide_images:
                    st.success(f"✅ {slide_count}枚のスライドをPNG画像に変換しました")
                    st.info(f"📋 変換方法: {'LibreOffice変換' if any('LibreOffice' in str(img) for img in slide_images) else '高品質フォールバック変換'}")
                else:
                    st.success(f"✅ {slide_count}枚のスライドを読み込みました")
                
                # スライド表示時間設定
                st.subheader("⏱️ スライド表示時間設定")
                
                if 'slide_durations' not in st.session_state:
                    st.session_state.slide_durations = [3.0] * slide_count
                
                st.write(f"**設定対象:** {slide_count}枚のスライド")
                
                # スライド毎の時間設定
                durations = []
                for i in range(slide_count):
                    col1, col2, col3 = st.columns([1, 2, 1])
                    with col1:
                        st.write(f"スライド {i+1}")
                    with col2:
                        duration = st.number_input(
                            f"表示時間（秒）",
                            min_value=0.5,
                            max_value=30.0,
                            value=st.session_state.slide_durations[i],
                            step=0.5,
                            key=f"duration_{i}"
                        )
                        durations.append(duration)
                        st.session_state.slide_durations[i] = duration
                    with col3:
                        # スライドプレビュー（簡易）
                        if i < len(slide_images) and os.path.exists(slide_images[i]):
                            try:
                                img = Image.open(slide_images[i])
                                # サムネイル作成
                                img.thumbnail((200, 150))
                                st.image(img, width=120)
                            except Exception as e:
                                st.text(f"プレビュー\nエラー: {str(e)[:20]}...")
                        else:
                            st.text(f"スライド {i+1}\n画像作成中...")
                
                total_duration = sum(durations)
                st.info(f"📊 総動画時間: {total_duration:.1f}秒")
                
                # 歌詞設定
                st.subheader("🎵 歌詞設定")
                add_lyrics = st.checkbox("歌詞を追加する")
                
                if add_lyrics:
                    if 'lyrics' not in st.session_state:
                        st.session_state.lyrics = []
                    
                    # 新しい歌詞を追加するフォーム
                    with st.expander("新しい歌詞を追加", expanded=True):
                        col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                        
                        with col1:
                            new_lyric_text = st.text_input(
                                "歌詞テキスト",
                                placeholder="歌詞を入力してください",
                                key="new_lyric_text"
                            )
                        with col2:
                            new_lyric_start = st.number_input(
                                "開始時間(秒)",
                                min_value=0.0,
                                max_value=total_duration,
                                value=0.0,
                                step=0.1,
                                key="new_lyric_start"
                            )
                        with col3:
                            new_lyric_end = st.number_input(
                                "終了時間(秒)",
                                min_value=0.1,
                                max_value=total_duration,
                                value=min(3.0, total_duration),
                                step=0.1,
                                key="new_lyric_end"
                            )
                        with col4:
                            st.write("") # スペース調整
                            if st.button("追加", type="primary", key="add_lyric"):
                                if new_lyric_text.strip():
                                    new_lyric = {
                                        "text": new_lyric_text,
                                        "start_time": new_lyric_start,
                                        "end_time": new_lyric_end
                                    }
                                    st.session_state.lyrics.append(new_lyric)
                                    st.rerun()
                    
                    # 既存の歌詞を表示・編集
                    if st.session_state.lyrics:
                        st.subheader("登録済み歌詞")
                        
                        for i, lyric in enumerate(st.session_state.lyrics):
                            with st.container():
                                col1, col2, col3, col4 = st.columns([3, 1, 1, 0.5])
                                
                                with col1:
                                    st.text_input(f"歌詞 {i+1}", value=lyric["text"], key=f"lyric_text_{i}", disabled=True)
                                with col2:
                                    st.text(f"{lyric['start_time']:.1f}秒")
                                with col3:
                                    st.text(f"{lyric['end_time']:.1f}秒")
                                with col4:
                                    if st.button("🗑️", key=f"delete_lyric_{i}", help="削除"):
                                        st.session_state.lyrics.pop(i)
                                        st.rerun()
                        
                        if st.button("全ての歌詞をクリア", type="secondary"):
                            st.session_state.lyrics = []
                            st.rerun()
                
                # 動画生成ボタン
                if st.button("🎬 スライド動画を生成", type="primary"):
                    if bgm_file is None:
                        st.warning("⚠️ BGMファイルがアップロードされていません")
                    
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    try:
                        # BGMファイルを一時保存
                        bgm_path = None
                        if bgm_file is not None:
                            bgm_path = os.path.join(temp_dir, "bgm.mp3")
                            bgm_file.seek(0)
                            with open(bgm_path, 'wb') as f:
                                f.write(bgm_file.read())
                        
                        status_text.text("スライド動画を生成中...")
                        progress_bar.progress(20)
                        
                        # 歌詞データの準備
                        lyrics_data = st.session_state.lyrics if add_lyrics else []
                        
                        # 動画出力パス
                        output_path = os.path.join(temp_dir, "slide_video.mp4")
                        
                        progress_bar.progress(50)
                        
                        # 動画を生成
                        create_slide_video(
                            slide_images,
                            durations,
                            bgm_path,
                            lyrics_data,
                            output_path
                        )
                        
                        progress_bar.progress(100)
                        status_text.text("生成完了！")
                        
                        # プレビュー表示
                        st.subheader("📹 プレビュー")
                        st.video(output_path, width=400)
                        
                        # ダウンロードボタンを表示
                        with open(output_path, 'rb') as file:
                            st.download_button(
                                label="📱 スライド動画をダウンロード",
                                data=file.read(),
                                file_name=f"slide_video_{pptx_file.name.replace('.pptx', '')}.mp4",
                                mime="video/mp4"
                            )
                        
                        st.success("✅ スライド動画の生成が完了しました！")
                        
                    except Exception as e:
                        st.error(f"❌ エラーが発生しました: {str(e)}")
        
        except Exception as e:
            st.error(f"❌ PowerPointファイルの処理に失敗しました: {str(e)}")

